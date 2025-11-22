#!/usr/bin/env python3
"""
Brain-Bridges PowerPoint Generator V3
With correctly configured Slide Masters for easy maintenance
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime
import os
import shutil
from PIL import Image, ImageDraw

# Import all design tokens (colors, fonts, layouts)
from design_tokens import *

def add_rounded_corners_to_image(image_path, corner_radius, border_color, border_width, output_path):
    """
    Add rounded corners and border to an image using PIL.

    Args:
        image_path: Path to input image
        corner_radius: Radius of corners in pixels
        border_color: RGB tuple for border (e.g., (77, 171, 247))
        border_width: Border width in pixels
        output_path: Path to save modified image
    """
    # Open image
    img = Image.open(image_path).convert("RGBA")

    # Create new image with space for border
    new_width = img.size[0] + 2 * border_width
    new_height = img.size[1] + 2 * border_width
    new_img = Image.new('RGBA', (new_width, new_height), (0, 0, 0, 0))

    # Paste original image in center
    new_img.paste(img, (border_width, border_width))

    # Create a mask with rounded corners for the entire new image
    mask = Image.new('L', (new_width, new_height), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle(
        [(0, 0), (new_width, new_height)],
        radius=corner_radius,
        fill=255
    )

    # Apply mask
    new_img.putalpha(mask)

    # Draw border
    draw = ImageDraw.Draw(new_img)
    draw.rounded_rectangle(
        [(0, 0), (new_width - 1, new_height - 1)],
        radius=corner_radius,
        outline=border_color,
        width=border_width
    )

    # Save with transparency
    new_img.save(output_path, "PNG")
    return output_path

def apply_master_elements(slide, slide_num, total_slides=17):
    """
    Applies master elements to a slide:
    - Background color
    - Logo "BRAIN BRIDGES" top left
    - Slide counter top right

    This function simulates a Slide Master, since python-pptx
    does not allow direct master editing.
    """

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BACKGROUND_DARK

    # Logo "BRAIN BRIDGES" top left (without "v: xii")
    logo_box = slide.shapes.add_textbox(
        LOGO_X, LOGO_Y,
        LOGO_WIDTH, LOGO_HEIGHT
    )
    logo_frame = logo_box.text_frame
    logo_frame.text = LOGO_TEXT
    logo_p = logo_frame.paragraphs[0]
    logo_p.font.size = FONT_SIZE_LOGO
    logo_p.font.bold = FONT_BOLD_LOGO
    logo_p.font.color.rgb = FONT_COLOR_LOGO
    # Letter-spacing
    for run in logo_p.runs:
        run.font.character_spacing = FONT_LETTER_SPACING_LOGO

    # Slide counter top right
    num_box = slide.shapes.add_textbox(
        SLIDE_NUMBER_X, SLIDE_NUMBER_Y,
        SLIDE_NUMBER_WIDTH, SLIDE_NUMBER_HEIGHT
    )
    num_frame = num_box.text_frame
    num_frame.text = f"{slide_num:02d}/{total_slides:02d}"
    num_p = num_frame.paragraphs[0]
    num_p.alignment = PP_ALIGN.RIGHT
    num_p.font.size = FONT_SIZE_SLIDE_NUMBER
    num_p.font.bold = FONT_BOLD_SLIDE_NUMBER
    num_p.font.color.rgb = FONT_COLOR_SLIDE_NUMBER
    
    return slide

def create_slide_1(prs):
    """Slide 1: THE AI PARADOX"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 1, 25)

    # The three keywords - using KEYWORD_THEME_PROBLEM
    keywords = [
        {"text": "THE", "color": KEYWORD_THEME_PROBLEM[0]},
        {"text": "AI", "color": KEYWORD_THEME_PROBLEM[1]},
        {"text": "PARADOX", "color": KEYWORD_THEME_PROBLEM[2]}
    ]

    for i, keyword in enumerate(keywords):
        y_pos = KEYWORD_Y_START + (i * KEYWORD_Y_GAP)

        keyword_box = slide.shapes.add_textbox(
            KEYWORD_BOX_X, Inches(y_pos),
            KEYWORD_BOX_WIDTH, KEYWORD_BOX_HEIGHT
        )
        tf = keyword_box.text_frame
        tf.text = keyword["text"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_KEYWORD
        p.font.bold = FONT_BOLD_KEYWORD
        p.font.color.rgb = keyword["color"]

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_KEYWORD  # Inter ExtraLight (font-weight: 200)
            run.font.character_spacing = FONT_LETTER_SPACING_KEYWORD

    return prs

def create_slide_2(prs):
    """Slide 2: Organisations want AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 2, 25)

    # Fixed header
    title_box = slide.shapes.add_textbox(
        CONTENT_HEADER_X, CONTENT_HEADER_Y,
        CONTENT_HEADER_WIDTH, CONTENT_HEADER_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "Organisations want AI"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CONTENT_TITLE
    p.font.bold = FONT_BOLD_CONTENT_TITLE
    p.font.color.rgb = FONT_COLOR_CONTENT_TITLE

    # CRITICAL: Font name must be set at RUN level!
    for run in p.runs:
        run.font.name = FONT_FAMILY_TITLE  # Inter ExtraLight (font-weight: 200)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        CONTENT_SUBTITLE_X, CONTENT_SUBTITLE_Y,
        CONTENT_SUBTITLE_WIDTH, CONTENT_SUBTITLE_HEIGHT
    )
    tf = subtitle_box.text_frame
    tf.text = "but can't have it ¯\\_(ツ)_/¯"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CONTENT_SUBTITLE
    p.font.bold = FONT_BOLD_CONTENT_SUBTITLE
    p.font.color.rgb = FONT_COLOR_CONTENT_SUBTITLE_ALERT

    # CRITICAL: Font name must be set at RUN level!
    for run in p.runs:
        run.font.name = FONT_FAMILY_SUBTITLE  # Menlo (monospace)

    # Problem items grid with PNG icons
    problems = [
        {
            "icon_key": "legal",
            "title": "Legal Practices",
            "desc": "Can't send client contracts to OpenAI",
            "violation": "ATTORNEY CLIENT PRIVILEGE"
        },
        {
            "icon_key": "medical",
            "title": "Medical Practices",
            "desc": "Can't upload patient records to ChatGPT",
            "violation": "HIPAA VIOLATIONS"
        },
        {
            "icon_key": "financial",
            "title": "Financial Services",
            "desc": "Can't process loan applications through Claude",
            "violation": "REGULATORY COMPLIANCE"
        },
        {
            "icon_key": "engineering",
            "title": "Engineering Teams",
            "desc": "Can't share R&D documents with AI",
            "violation": "TRADE SECRETS"
        }
    ]

    for i, problem in enumerate(problems):
        # 2x2 Grid Layout: Calculate row and column
        row = i // 2  # 0 or 1 (top row or bottom row)
        col = i % 2   # 0 or 1 (left column or right column)

        x = Inches(PROBLEM_GRID_X_POSITIONS[col])
        y_start = Inches(PROBLEM_GRID_Y_POSITIONS[row])
        box_width = Inches(PROBLEM_GRID_BOX_WIDTH)
        box_height = Inches(PROBLEM_GRID_BOX_HEIGHT)

        # Card background with rounded corners and border
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y_start,
            box_width, box_height
        )
        # Fill
        card.fill.solid()
        card.fill.fore_color.rgb = PROBLEM_CARD_FILL_COLOR
        # Border
        card.line.color.rgb = PROBLEM_CARD_BORDER_COLOR
        card.line.width = PROBLEM_CARD_BORDER_WIDTH
        # Rounded corners
        card.adjustments[0] = 0.05  # Corner radius adjustment

        # Icon (PNG image)
        icon_x = x + Inches(PROBLEM_ICON_X_OFFSET)
        icon_y = y_start + Inches(PROBLEM_ICON_Y_OFFSET)
        icon_path = PROBLEM_ICONS[problem["icon_key"]]

        slide.shapes.add_picture(
            icon_path,
            icon_x, icon_y,
            width=PROBLEM_ICON_WIDTH
        )

        # Title
        title_box = slide.shapes.add_textbox(
            x, y_start + Inches(PROBLEM_TITLE_Y_OFFSET),
            box_width, Inches(PROBLEM_TITLE_HEIGHT)
        )
        tf = title_box.text_frame
        tf.text = problem["title"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_PROBLEM_TITLE
        p.font.bold = FONT_BOLD_PROBLEM_TITLE
        p.font.color.rgb = FONT_COLOR_PROBLEM_TITLE

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_TITLE  # Inter ExtraLight (font-weight: 200)

        # Description
        desc_box = slide.shapes.add_textbox(
            x, y_start + Inches(PROBLEM_DESC_Y_OFFSET),
            box_width, Inches(PROBLEM_DESC_HEIGHT)
        )
        tf = desc_box.text_frame
        tf.text = problem["desc"]
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_PROBLEM_DESC
        p.font.bold = FONT_BOLD_PROBLEM_DESC
        p.font.color.rgb = FONT_COLOR_PROBLEM_DESC

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_TITLE  # Same as title (Inter ExtraLight)

        # Violation
        viol_box = slide.shapes.add_textbox(
            x, y_start + Inches(PROBLEM_VIOLATION_Y_OFFSET),
            box_width, Inches(PROBLEM_VIOLATION_HEIGHT)
        )
        tf = viol_box.text_frame
        tf.text = problem["violation"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_PROBLEM_VIOLATION
        p.font.bold = FONT_BOLD_PROBLEM_VIOLATION
        p.font.color.rgb = FONT_COLOR_PROBLEM_VIOLATION

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_VIOLATION  # Menlo (monospace)

    return prs

def create_slide_3(prs):
    """Slide 3: Market Reality"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 3, 25)

    # Fixed header
    title_box = slide.shapes.add_textbox(
        CONTENT_HEADER_X, CONTENT_HEADER_Y,
        CONTENT_HEADER_WIDTH, CONTENT_HEADER_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "Market Reality"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CONTENT_TITLE
    p.font.bold = FONT_BOLD_CONTENT_TITLE
    p.font.color.rgb = FONT_COLOR_CONTENT_TITLE

    # CRITICAL: Font name must be set at RUN level!
    for run in p.runs:
        run.font.name = FONT_FAMILY_TITLE  # Inter ExtraLight (font-weight: 200)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        CONTENT_SUBTITLE_X, CONTENT_SUBTITLE_Y,
        CONTENT_SUBTITLE_WIDTH, CONTENT_SUBTITLE_HEIGHT
    )
    tf = subtitle_box.text_frame
    tf.text = "Massive demand blocked by fundamental constraints"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CONTENT_SUBTITLE
    p.font.bold = FONT_BOLD_CONTENT_SUBTITLE
    p.font.color.rgb = FONT_COLOR_CONTENT_SUBTITLE_ALERT  # Red like Slide 2

    # CRITICAL: Font name must be set at RUN level!
    for run in p.runs:
        run.font.name = FONT_FAMILY_SUBTITLE  # Menlo monospace like Slide 2

    # Large stat: $1.7T
    large_stat_box = slide.shapes.add_textbox(
        LARGE_STAT_X, LARGE_STAT_Y,
        LARGE_STAT_WIDTH, LARGE_STAT_HEIGHT
    )
    tf = large_stat_box.text_frame
    tf.text = "$1.7T"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_LARGE_STAT_NUMBER
    p.font.bold = False
    p.font.color.rgb = COLOR_ACCENT_CYAN

    # CRITICAL: Font name must be set at RUN level!
    for run in p.runs:
        run.font.name = FONT_FAMILY_STAT_NUMBER  # Inter ExtraLight (font-weight: 200)

    # Large stat label
    large_stat_label_box = slide.shapes.add_textbox(
        LARGE_STAT_X, LARGE_STAT_LABEL_Y,
        LARGE_STAT_WIDTH, Inches(0.4)
    )
    tf = large_stat_label_box.text_frame
    tf.text = "Global AI market by 2032"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_LARGE_STAT_LABEL
    p.font.bold = False
    p.font.color.rgb = COLOR_TEXT_WHITE

    # CRITICAL: Font name must be set at RUN level!
    for run in p.runs:
        run.font.name = FONT_FAMILY_STAT_LABEL  # Inter Light (font-weight: 300)

    # Stat cards (4 cards in a row)
    stats = [
        {
            "number": "96%",
            "label": "Want AI expansion",
            "source": "McKinsey Global AI Survey 2024"
        },
        {
            "number": "53%",
            "label": "Blocked by data privacy",
            "source": "Deloitte Enterprise AI Study 2024"
        },
        {
            "number": "30%",
            "label": "Projects abandoned after POC",
            "source": "MIT Technology Review 2024"
        },
        {
            "number": "40%",
            "label": "Healthcare AI apps blocked",
            "source": "HIMSS Healthcare IT Report 2024"
        }
    ]

    for i, stat in enumerate(stats):
        x = Inches(STAT_CARD_X_POSITIONS[i])
        y_start = Inches(STAT_CARD_Y_START)
        card_width = Inches(STAT_CARD_WIDTH)
        card_height = Inches(STAT_CARD_HEIGHT)

        # Card background with rounded corners
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y_start,
            card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = STAT_CARD_FILL_COLOR
        card.line.color.rgb = STAT_CARD_BORDER_COLOR
        card.line.width = STAT_CARD_BORDER_WIDTH
        card.adjustments[0] = 0.05

        # Stat number
        number_box = slide.shapes.add_textbox(
            x, y_start + Inches(STAT_NUMBER_Y_OFFSET),
            card_width, Inches(STAT_NUMBER_HEIGHT)
        )
        tf = number_box.text_frame
        tf.text = stat["number"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_STAT_NUMBER
        p.font.bold = False
        p.font.color.rgb = COLOR_ACCENT_CYAN

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_STAT_NUMBER  # Inter ExtraLight (font-weight: 200)

        # Stat label
        label_box = slide.shapes.add_textbox(
            x, y_start + Inches(STAT_LABEL_Y_OFFSET),
            card_width, Inches(STAT_LABEL_HEIGHT)
        )
        tf = label_box.text_frame
        tf.text = stat["label"]
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_STAT_LABEL
        p.font.bold = False
        p.font.color.rgb = COLOR_TEXT_WHITE

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_STAT_LABEL  # Inter Light (font-weight: 300)

        # Stat source
        source_box = slide.shapes.add_textbox(
            x, y_start + Inches(STAT_SOURCE_Y_OFFSET),
            card_width, Inches(STAT_SOURCE_HEIGHT)
        )
        tf = source_box.text_frame
        tf.text = stat["source"]
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_STAT_SOURCE
        p.font.bold = False
        p.font.color.rgb = FONT_COLOR_STAT_SOURCE

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_PRIMARY  # Inter Regular (small body text)

    return prs

def create_slide_4(prs):
    """Slide 4: SOVEREIGN AI SOLUTION"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 4, 25)

    # The three keywords - using KEYWORD_THEME_SOLUTION
    keywords = [
        {"text": "SOVEREIGN", "color": KEYWORD_THEME_SOLUTION[0]},
        {"text": "AI", "color": KEYWORD_THEME_SOLUTION[1]},
        {"text": "SOLUTION", "color": KEYWORD_THEME_SOLUTION[2]}
    ]

    for i, keyword in enumerate(keywords):
        y_pos = KEYWORD_Y_START + (i * KEYWORD_Y_GAP)

        keyword_box = slide.shapes.add_textbox(
            KEYWORD_BOX_X, Inches(y_pos),
            KEYWORD_BOX_WIDTH, KEYWORD_BOX_HEIGHT
        )
        tf = keyword_box.text_frame
        tf.text = keyword["text"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_KEYWORD
        p.font.bold = FONT_BOLD_KEYWORD
        p.font.color.rgb = keyword["color"]

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_KEYWORD  # Inter ExtraLight (font-weight: 200)
            run.font.character_spacing = FONT_LETTER_SPACING_KEYWORD

    return prs

def create_slide_5(prs):
    """Slide 5: BRAIN-BRIDGES Introduction (like Slide 6 but with text instead of features)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 5, 25)

    # =========================================================================
    # LEFT SIDE: Title, Subtitle, Description Text
    # =========================================================================

    # Hero Title: "BRAIN-BRIDGES"
    title_box = slide.shapes.add_textbox(
        HERO_TITLE_X, HERO_TITLE_Y,
        HERO_TITLE_WIDTH, HERO_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "BRAIN-BRIDGES"
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.name = FONT_FAMILY_HERO_TITLE
    p.font.size = FONT_SIZE_HERO_TITLE
    p.font.bold = FONT_BOLD_HERO_TITLE
    p.font.color.rgb = COLOR_ACCENT_BLUE

    for run in p.runs:
        run.font.character_spacing = FONT_LETTER_SPACING_HERO_TITLE

    # Hero Subtitle: "SOVEREIGN AI FOR ORGANISATIONS"
    subtitle_box = slide.shapes.add_textbox(
        HERO_SUBTITLE_X, HERO_SUBTITLE_Y,
        HERO_SUBTITLE_WIDTH, HERO_SUBTITLE_HEIGHT
    )
    tf = subtitle_box.text_frame
    tf.text = "SOVEREIGN AI FOR ORGANISATIONS"
    p = tf.paragraphs[0]
    p.font.name = FONT_FAMILY_HERO_SUBTITLE
    p.font.size = FONT_SIZE_HERO_SUBTITLE
    p.font.bold = FONT_BOLD_HERO_SUBTITLE
    p.font.color.rgb = FONT_COLOR_HERO_SUBTITLE

    # Description text (3 paragraphs) - positioned below subtitle
    text_box = slide.shapes.add_textbox(
        HERO_FEATURES_X, Inches(3.5),
        HERO_FEATURES_WIDTH, Inches(4.0)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    # Paragraph 1
    p1 = tf.paragraphs[0]
    p1.text = "We offer an AI Bot that enables users to have chat-like conversations about their organization. The bot has access to the organization's documents repository."
    p1.font.size = Pt(14)
    p1.font.color.rgb = COLOR_TEXT_WHITE
    p1.space_after = Pt(16)
    for run in p1.runs:
        run.font.name = FONT_FAMILY_PRIMARY

    # Paragraph 2 (with bold words)
    p2 = tf.add_paragraph()
    p2.text = "To guarantee absolute data sovereignty, privacy, and security, we design the entire system – including AI model, database and documents – for full on-premises deployment."
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_after = Pt(16)

    # Make specific words bold
    text2 = p2.text
    p2.clear()

    parts = [
        ("To guarantee absolute ", False),
        ("data sovereignty", True),
        (", ", False),
        ("privacy", True),
        (", and ", False),
        ("security", True),
        (", we design the entire system – including AI model, database and documents – for full on-premises deployment.", False)
    ]

    for text, is_bold in parts:
        run = p2.add_run()
        run.text = text
        run.font.name = FONT_FAMILY_PRIMARY
        run.font.size = Pt(14)
        run.font.bold = is_bold
        run.font.color.rgb = COLOR_TEXT_WHITE

    # Paragraph 3
    p3 = tf.add_paragraph()
    p3.text = "Everything is delivered as a compact, ready-to-use system – just plug in and start."
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_WHITE
    for run in p3.runs:
        run.font.name = FONT_FAMILY_PRIMARY

    # =========================================================================
    # RIGHT SIDE: Product Image with Status Badge (NO TECH SPECS)
    # =========================================================================

    # Get actual image dimensions for aspect ratio
    img_actual_height = HERO_IMAGE_HEIGHT

    try:
        pil_img = Image.open(HERO_IMAGE_PATH)
        img_width_px, img_height_px = pil_img.size
        aspect_ratio = img_height_px / img_width_px
        img_actual_height = HERO_IMAGE_WIDTH * aspect_ratio
    except Exception as e:
        print(f"⚠️  Warning: Could not get image dimensions: {e}")

    # Process image: add rounded corners and border
    try:
        os.makedirs("temp", exist_ok=True)
        border_width_px = int(HERO_IMAGE_BORDER_WIDTH.pt * 1.33)
        rounded_image_path = "temp/hero_image_rounded.png"
        add_rounded_corners_to_image(
            HERO_IMAGE_PATH,
            HERO_IMAGE_CORNER_RADIUS_PX,
            HERO_IMAGE_BORDER_COLOR_RGB,
            border_width_px,
            rounded_image_path
        )

        product_img = slide.shapes.add_picture(
            rounded_image_path,
            HERO_IMAGE_X, HERO_IMAGE_Y,
            width=HERO_IMAGE_WIDTH
        )
        img_actual_height = product_img.height
    except Exception as e:
        print(f"⚠️  Warning: Could not load/process hero image: {e}")
        try:
            product_img = slide.shapes.add_picture(
                HERO_IMAGE_PATH,
                HERO_IMAGE_X, HERO_IMAGE_Y,
                width=HERO_IMAGE_WIDTH
            )
            img_actual_height = product_img.height
        except:
            pass

    # NO Status Badge on Slide 5 (only image)

    return prs

def create_slide_6(prs):
    """Slide 6: BRAIN-BRIDGES Hero Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 6, 25)

    # =========================================================================
    # LEFT SIDE: Title, Subtitle, Features
    # =========================================================================

    # Hero Title: "BRAIN-BRIDGES"
    # Note: PowerPoint gradient text is complex, using solid blue color
    title_box = slide.shapes.add_textbox(
        HERO_TITLE_X, HERO_TITLE_Y,
        HERO_TITLE_WIDTH, HERO_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "BRAIN-BRIDGES"
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.name = FONT_FAMILY_HERO_TITLE  # Inter-ExtraBold (font-weight: 800)
    p.font.size = FONT_SIZE_HERO_TITLE
    p.font.bold = FONT_BOLD_HERO_TITLE
    p.font.color.rgb = COLOR_ACCENT_BLUE  # Solid blue (gradient not reliable in pptx)

    # Apply negative letter-spacing (like HTML version)
    for run in p.runs:
        run.font.character_spacing = FONT_LETTER_SPACING_HERO_TITLE

    # Hero Subtitle: "SOVEREIGN AI FOR ORGANISATIONS"
    subtitle_box = slide.shapes.add_textbox(
        HERO_SUBTITLE_X, HERO_SUBTITLE_Y,
        HERO_SUBTITLE_WIDTH, HERO_SUBTITLE_HEIGHT
    )
    tf = subtitle_box.text_frame
    tf.text = "SOVEREIGN AI FOR ORGANISATIONS"
    p = tf.paragraphs[0]
    p.font.name = FONT_FAMILY_HERO_SUBTITLE  # Menlo (monospace font)
    p.font.size = FONT_SIZE_HERO_SUBTITLE
    p.font.bold = FONT_BOLD_HERO_SUBTITLE
    p.font.color.rgb = FONT_COLOR_HERO_SUBTITLE

    # Hero Features List (6 items - 5 with checkmarks, last with plug icon)
    features = [
        ("Multi-Model Support (Mistral, Llama, Gemma, others)", "check"),
        ("PostgreSQL with pgvector (Enterprise-Ready)", "check"),
        ("Production Ready Webinterface", "check"),
        ("Nomic Embeddings (Multi-Language)", "check"),
        ("MCP-Compatible Agent Framework", "check"),
        ("Zero-Config Deployment", "plug")  # Plug icon for last item
    ]

    for i, (feature_text, icon_type) in enumerate(features):
        y_pos = Inches(HERO_FEATURES_Y_START + i * HERO_FEATURE_GAP)

        # Add underline extending to image (drawn FIRST, behind text)
        # Made thinner and more subtle
        line_y = y_pos + Inches(HERO_FEATURE_HEIGHT) - Inches(0.02)
        line_start_x = HERO_FEATURES_X
        line_end_x = HERO_IMAGE_X + Inches(0.2)  # Extend a bit into image area

        underline = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR_TYPE.STRAIGHT
            line_start_x, line_y,
            line_end_x, line_y
        )
        # Thinner and more washed out
        underline.line.color.rgb = RGBColor(77, 171, 247)  # Blue
        underline.line.width = Pt(0.75)  # Reduced from 1.5
        underline.line.fill.solid()
        underline.line.fill.fore_color.rgb = RGBColor(77, 171, 247)
        underline.line.fill.transparency = 0.85  # More transparent (85% vs 80%)

        # Icon (checkmark or plug)
        try:
            icon_x = HERO_FEATURES_X
            icon_y_center = y_pos + Inches(HERO_FEATURE_HEIGHT / 2)
            icon_y = icon_y_center - (HERO_FEATURE_ICON_SIZE / 2)

            # Use plug icon for last item, checkmark for others
            icon_path = HERO_STATUS_ICON if icon_type == "plug" else HERO_FEATURE_CHECKMARK_ICON

            icon = slide.shapes.add_picture(
                icon_path,
                icon_x, icon_y,
                width=HERO_FEATURE_ICON_SIZE
            )
        except Exception as e:
            print(f"⚠️  Warning: Could not load icon: {e}")

        # Feature text box (next to icon)
        text_x = HERO_FEATURES_X + HERO_FEATURE_ICON_SIZE + Inches(0.15)  # Gap after icon
        text_width = HERO_FEATURES_WIDTH - HERO_FEATURE_ICON_SIZE - Inches(0.15)

        feature_box = slide.shapes.add_textbox(
            text_x, y_pos,
            text_width, Inches(HERO_FEATURE_HEIGHT)
        )
        tf = feature_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text to align with icon

        # Feature text (white)
        p = tf.paragraphs[0]
        p.font.name = FONT_FAMILY_HERO_FEATURE  # Inter-Light (font-weight: 300)
        p.font.size = FONT_SIZE_HERO_FEATURE
        p.font.color.rgb = FONT_COLOR_HERO_FEATURE
        p.text = feature_text

    # =========================================================================
    # RIGHT SIDE: Product Image with border, Status Badge, Tech Specs
    # =========================================================================

    # Get actual image dimensions for aspect ratio
    img_actual_height = HERO_IMAGE_HEIGHT

    try:
        pil_img = Image.open(HERO_IMAGE_PATH)
        img_width_px, img_height_px = pil_img.size
        aspect_ratio = img_height_px / img_width_px
        img_actual_height = HERO_IMAGE_WIDTH * aspect_ratio
    except Exception as e:
        print(f"⚠️  Warning: Could not get image dimensions: {e}")

    # Process image: add rounded corners and border
    try:
        # Create temp directory if it doesn't exist
        os.makedirs("temp", exist_ok=True)

        # Convert border width from Pt to pixels (approximate: 1pt ≈ 1.33px)
        border_width_px = int(HERO_IMAGE_BORDER_WIDTH.pt * 1.33)

        # Create rounded corner version of image with border
        rounded_image_path = "temp/hero_image_rounded.png"
        add_rounded_corners_to_image(
            HERO_IMAGE_PATH,
            HERO_IMAGE_CORNER_RADIUS_PX,
            HERO_IMAGE_BORDER_COLOR_RGB,
            border_width_px,
            rounded_image_path
        )

        # Add image with rounded corners and border
        product_img = slide.shapes.add_picture(
            rounded_image_path,
            HERO_IMAGE_X, HERO_IMAGE_Y,
            width=HERO_IMAGE_WIDTH
        )

        # Get actual image dimensions after adding
        img_actual_height = product_img.height
    except Exception as e:
        print(f"⚠️  Warning: Could not load/process hero image: {e}")
        # Fallback: use original image without rounded corners
        try:
            product_img = slide.shapes.add_picture(
                HERO_IMAGE_PATH,
                HERO_IMAGE_X, HERO_IMAGE_Y,
                width=HERO_IMAGE_WIDTH
            )
            img_actual_height = product_img.height
        except:
            pass

    # Calculate positions for cards INSIDE the image
    img_bottom = HERO_IMAGE_Y + img_actual_height
    img_right = HERO_IMAGE_X + HERO_IMAGE_WIDTH

    # Status Badge: INSIDE image at top right with margin
    status_x = img_right - HERO_STATUS_WIDTH - HERO_STATUS_MARGIN
    status_y = HERO_IMAGE_Y + HERO_STATUS_MARGIN

    status_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        status_x, status_y,
        HERO_STATUS_WIDTH, HERO_STATUS_HEIGHT
    )
    # Semi-transparent effect using slightly lighter color
    fill = status_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 32, 45)  # Lighter dark background for semi-transparent effect
    status_shape.line.color.rgb = COLOR_ACCENT_CYAN
    status_shape.line.width = Pt(1)

    # Plug icon and text - CENTERED in badge
    # Icon left side, text right side, both vertically centered
    badge_center_y = status_y + (HERO_STATUS_HEIGHT / 2)

    try:
        # Icon (left-center of badge)
        icon_x = status_x + Inches(0.15)
        icon_y = badge_center_y - (HERO_STATUS_ICON_SIZE / 2)

        plug_icon = slide.shapes.add_picture(
            HERO_STATUS_ICON,
            icon_x, icon_y,
            width=HERO_STATUS_ICON_SIZE
        )
    except Exception as e:
        print(f"⚠️  Warning: Could not load plug icon: {e}")

    # Text (right of icon, centered)
    text_x = status_x + Inches(0.33)
    text_box = slide.shapes.add_textbox(
        text_x, status_y,
        HERO_STATUS_WIDTH - Inches(0.33), HERO_STATUS_HEIGHT
    )
    tf = text_box.text_frame
    tf.text = HERO_STATUS_TEXT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text to align with icon
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_FAMILY_HERO_STATUS  # Inter-Medium (font-weight: 500)
    p.font.size = FONT_SIZE_HERO_STATUS
    p.font.bold = True
    p.font.color.rgb = FONT_COLOR_HERO_STATUS

    # Tech Specs (3 columns: Processor, Memory, Users) INSIDE image at bottom
    spec_labels = ["PROCESSOR", "MEMORY", "USERS"]
    spec_values = [
        HERO_SPECS_CONFIG["processor"],
        HERO_SPECS_CONFIG["memory"],
        HERO_SPECS_CONFIG["users"]
    ]

    # Calculate card width: (image_width - 2*margin - 2*gap) / 3
    total_card_width = HERO_IMAGE_WIDTH - (2 * HERO_SPECS_MARGIN)
    single_card_width = (total_card_width - (2 * HERO_SPECS_GAP)) / 3

    # Cards at bottom of image with margin
    specs_y = img_bottom - HERO_SPECS_HEIGHT - HERO_SPECS_MARGIN

    for i, (label, value) in enumerate(zip(spec_labels, spec_values)):
        x_pos = HERO_IMAGE_X + HERO_SPECS_MARGIN + i * (single_card_width + HERO_SPECS_GAP)

        # Spec card background with rounded corners - SEMI-TRANSPARENT effect
        spec_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, specs_y,
            single_card_width, HERO_SPECS_HEIGHT
        )
        # Semi-transparent effect using slightly lighter color (transparency doesn't work reliably)
        fill = spec_card.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(25, 32, 45)  # Lighter dark background for semi-transparent effect

        # Blue border
        spec_card.line.color.rgb = COLOR_ACCENT_BLUE
        spec_card.line.width = Pt(0.75)

        # Rounded corners
        spec_card.adjustments[0] = 0.05

        # Label (top)
        label_box = slide.shapes.add_textbox(
            x_pos, specs_y + Inches(0.18),
            single_card_width, Inches(0.25)
        )
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_HERO_SPEC_LABEL
        p.font.color.rgb = FONT_COLOR_HERO_SPEC_LABEL
        p.font.bold = True

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_HERO_SPEC_LABEL  # Inter-SemiBold (font-weight: 600)

        # Value (bottom)
        value_box = slide.shapes.add_textbox(
            x_pos, specs_y + Inches(0.48),
            single_card_width, Inches(0.4)
        )
        tf = value_box.text_frame
        tf.text = value
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_HERO_SPEC_VALUE
        p.font.color.rgb = FONT_COLOR_HERO_SPEC_VALUE
        p.font.bold = FONT_BOLD_HERO_SPEC_VALUE

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_HERO_SPEC_VALUE  # Inter-Bold (font-weight: 700)

    return prs

def create_slide_7(prs):
    """Slide 7: UNDERSTANDING INFERENCE MECHANICS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 7, 25)

    # The three keywords - using KEYWORD_THEME_TECH
    keywords = [
        {"text": "UNDERSTANDING", "color": KEYWORD_THEME_TECH[0]},
        {"text": "INFERENCE", "color": KEYWORD_THEME_TECH[1]},
        {"text": "MECHANICS", "color": KEYWORD_THEME_TECH[2]}
    ]

    for i, keyword in enumerate(keywords):
        y_pos = KEYWORD_Y_START + (i * KEYWORD_Y_GAP)

        keyword_box = slide.shapes.add_textbox(
            KEYWORD_BOX_X, Inches(y_pos),
            KEYWORD_BOX_WIDTH, KEYWORD_BOX_HEIGHT
        )
        tf = keyword_box.text_frame
        tf.text = keyword["text"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_KEYWORD
        p.font.bold = FONT_BOLD_KEYWORD
        p.font.color.rgb = keyword["color"]

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_KEYWORD  # Inter ExtraLight (font-weight: 200)
            run.font.character_spacing = FONT_LETTER_SPACING_KEYWORD

    return prs

def create_slide_8(prs):
    """Slide 8: Tokenization Intro - A Sample from legal domain"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 8, 25)

    # Title: "A Sample from legal domain:"
    title_box = slide.shapes.add_textbox(
        TOKENIZATION_TITLE_X, TOKENIZATION_TITLE_Y,
        TOKENIZATION_TITLE_WIDTH, TOKENIZATION_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "A Sample from legal domain:"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_TOKENIZATION_TITLE
    p.font.color.rgb = FONT_COLOR_TOKENIZATION_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_TOKENIZATION_TITLE

    # Arrow down
    arrow_box = slide.shapes.add_textbox(
        TOKENIZATION_ARROW_X, TOKENIZATION_ARROW_Y,
        TOKENIZATION_ARROW_WIDTH, TOKENIZATION_ARROW_HEIGHT
    )
    tf = arrow_box.text_frame
    tf.text = TOKENIZATION_ARROW_TEXT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_TOKENIZATION_ARROW
    p.font.color.rgb = FONT_COLOR_TOKENIZATION_ARROW

    # Token boxes in horizontal row
    for i, token_text in enumerate(TOKENIZATION_TOKENS):
        x_pos = TOKENIZATION_TOKEN_X_START + (i * (TOKENIZATION_TOKEN_WIDTH + TOKENIZATION_TOKEN_GAP))

        # Create token box
        token_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, TOKENIZATION_TOKENS_Y,
            TOKENIZATION_TOKEN_WIDTH, TOKENIZATION_TOKEN_HEIGHT
        )
        token_box.fill.solid()
        token_box.fill.fore_color.rgb = TOKENIZATION_TOKEN_FILL_COLOR
        token_box.line.color.rgb = TOKENIZATION_TOKEN_BORDER_COLOR
        token_box.line.width = TOKENIZATION_TOKEN_BORDER_WIDTH

        # Token text
        tf = token_box.text_frame
        tf.text = token_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_TOKENIZATION_TOKEN
        p.font.color.rgb = FONT_COLOR_TOKENIZATION_TOKEN
        for run in p.runs:
            run.font.name = FONT_FAMILY_TOKENIZATION_TOKEN

    return prs

def create_slide_9(prs):
    """Slide 9: Vector Embeddings (Token → Vector Lookup)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 9, 25)

    # Create each token row (Wit, nesses, must, tell, nothing)
    for i, token_info in enumerate(TOKEN_DATA):
        y_pos = TOKEN_ROW_Y_START + (i * TOKEN_ROW_GAP)

        # Token box (left side)
        token_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            TOKEN_BOX_X, Inches(y_pos),
            TOKEN_BOX_WIDTH, TOKEN_BOX_HEIGHT
        )
        token_box.fill.solid()
        token_box.fill.fore_color.rgb = TOKEN_BOX_FILL_COLOR
        token_box.line.color.rgb = TOKEN_BOX_BORDER_COLOR
        token_box.line.width = TOKEN_BOX_BORDER_WIDTH

        # Token text
        tf = token_box.text_frame
        tf.text = token_info["token"]
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_TOKEN
        p.font.color.rgb = FONT_COLOR_TOKEN
        for run in p.runs:
            run.font.name = FONT_FAMILY_TOKEN

        # Arrow (center)
        arrow_box = slide.shapes.add_textbox(
            ARROW_X, Inches(y_pos),
            ARROW_WIDTH, TOKEN_BOX_HEIGHT
        )
        tf = arrow_box.text_frame
        tf.text = ARROW_TEXT
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_ARROW
        p.font.color.rgb = FONT_COLOR_ARROW

        # Vector cells (right side, 6 cells)
        for j, vector_value in enumerate(token_info["vectors"]):
            cell_x = VECTOR_GRID_X + (j * (VECTOR_CELL_WIDTH + VECTOR_CELL_GAP))

            # Create vector cell box
            vector_cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                cell_x, Inches(y_pos),
                VECTOR_CELL_WIDTH, TOKEN_BOX_HEIGHT
            )
            vector_cell.fill.solid()
            vector_cell.fill.fore_color.rgb = VECTOR_CELL_FILL_COLOR
            vector_cell.line.color.rgb = VECTOR_CELL_BORDER_COLOR
            vector_cell.line.width = VECTOR_CELL_BORDER_WIDTH

            # Vector value text
            tf = vector_cell.text_frame
            tf.text = vector_value
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = FONT_SIZE_VECTOR
            p.font.color.rgb = FONT_COLOR_VECTOR

            # Last cell ("...") should be accent color
            if vector_value == "...":
                p.font.color.rgb = COLOR_ACCENT_BLUE

            for run in p.runs:
                run.font.name = FONT_FAMILY_VECTOR

    return prs

def create_slide_10(prs):
    """Slide 10: Attention is all you need - Attention Matrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 10, 25)

    # Title
    title_box = slide.shapes.add_textbox(
        ATTENTION_TITLE_X, ATTENTION_TITLE_Y,
        ATTENTION_TITLE_WIDTH, ATTENTION_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "Attention is all you need"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ATTENTION_TITLE
    p.font.color.rgb = FONT_COLOR_ATTENTION_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_ATTENTION_TITLE

    # Helper function to get score color
    def get_score_style(score):
        """Returns (fill_color, border_color, text_color) for a given score"""
        if score >= 0.8:
            return (ATTENTION_SCORE_HIGH_FILL, ATTENTION_SCORE_HIGH_BORDER, ATTENTION_SCORE_HIGH_COLOR)
        elif score >= 0.05:
            return (ATTENTION_SCORE_MED_FILL, ATTENTION_SCORE_MED_BORDER, ATTENTION_SCORE_MED_COLOR)
        elif score >= 0.02:
            return (ATTENTION_SCORE_LOW_MED_FILL, ATTENTION_SCORE_LOW_MED_BORDER, ATTENTION_SCORE_LOW_MED_COLOR)
        else:
            return (ATTENTION_SCORE_LOW_FILL, ATTENTION_SCORE_LOW_BORDER, ATTENTION_SCORE_LOW_COLOR)

    # Create 6x6 grid (header row + 5 data rows, header col + 5 data cols)
    # Row 0: Header row
    # Col 0: Empty top-left cell
    # Cols 1-5: Token headers (Wit, nesses, must, tell, nothing)

    # Empty top-left cell (transparent, no content)
    # (We skip creating it since it's just empty space)

    # Header row: Token headers (columns)
    for col_idx, token in enumerate(ATTENTION_TOKENS):
        x_pos = ATTENTION_MATRIX_X + ((col_idx + 1) * (ATTENTION_CELL_WIDTH + ATTENTION_CELL_GAP))
        y_pos = ATTENTION_MATRIX_Y

        cell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos,
            ATTENTION_CELL_WIDTH, ATTENTION_CELL_HEIGHT
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = ATTENTION_HEADER_FILL_COLOR
        cell.line.color.rgb = ATTENTION_HEADER_BORDER_COLOR
        cell.line.width = ATTENTION_CELL_BORDER_WIDTH

        tf = cell.text_frame
        tf.text = token
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_ATTENTION_HEADER
        p.font.color.rgb = FONT_COLOR_ATTENTION_HEADER
        for run in p.runs:
            run.font.name = FONT_FAMILY_ATTENTION_HEADER

    # Data rows (rows 1-5)
    for row_idx, (row_token, row_scores) in enumerate(zip(ATTENTION_TOKENS, ATTENTION_MATRIX_DATA)):
        y_pos = ATTENTION_MATRIX_Y + ((row_idx + 1) * (ATTENTION_CELL_HEIGHT + ATTENTION_CELL_GAP))

        # Row header (token name)
        x_pos = ATTENTION_MATRIX_X
        cell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos,
            ATTENTION_CELL_WIDTH, ATTENTION_CELL_HEIGHT
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = ATTENTION_HEADER_FILL_COLOR
        cell.line.color.rgb = ATTENTION_HEADER_BORDER_COLOR
        cell.line.width = ATTENTION_CELL_BORDER_WIDTH

        tf = cell.text_frame
        tf.text = row_token
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_ATTENTION_HEADER
        p.font.color.rgb = FONT_COLOR_ATTENTION_HEADER
        for run in p.runs:
            run.font.name = FONT_FAMILY_ATTENTION_HEADER

        # Score cells
        for col_idx, score in enumerate(row_scores):
            x_pos = ATTENTION_MATRIX_X + ((col_idx + 1) * (ATTENTION_CELL_WIDTH + ATTENTION_CELL_GAP))

            fill_color, border_color, text_color = get_score_style(score)

            cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, y_pos,
                ATTENTION_CELL_WIDTH, ATTENTION_CELL_HEIGHT
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            cell.line.color.rgb = border_color
            cell.line.width = ATTENTION_CELL_BORDER_WIDTH

            tf = cell.text_frame
            tf.text = f"{score:.2f}"
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = FONT_SIZE_ATTENTION_SCORE
            p.font.color.rgb = text_color
            for run in p.runs:
                run.font.name = FONT_FAMILY_ATTENTION_SCORE

    # Footnote
    footnote_box = slide.shapes.add_textbox(
        ATTENTION_FOOTNOTE_X, ATTENTION_FOOTNOTE_Y,
        ATTENTION_FOOTNOTE_WIDTH, ATTENTION_FOOTNOTE_HEIGHT
    )
    tf = footnote_box.text_frame
    tf.text = ATTENTION_FOOTNOTE_TEXT
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = FONT_SIZE_ATTENTION_FOOTNOTE
    p.font.color.rgb = FONT_COLOR_ATTENTION_FOOTNOTE
    p.font.italic = True

    return prs

def create_slide_11(prs):
    """Slide 11: Next word prediction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 11, 25)

    # Title
    title_box = slide.shapes.add_textbox(
        PREDICTION_TITLE_X, PREDICTION_TITLE_Y,
        PREDICTION_TITLE_WIDTH, PREDICTION_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "Next word prediction"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_PREDICTION_TITLE
    p.font.color.rgb = FONT_COLOR_PREDICTION_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_PREDICTION_TITLE

    # Context Vector Bar
    vector_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PREDICTION_CONTENT_X, PREDICTION_VECTOR_Y,
        PREDICTION_CONTENT_WIDTH, PREDICTION_VECTOR_HEIGHT
    )
    vector_bar.fill.solid()
    vector_bar.fill.fore_color.rgb = PREDICTION_VECTOR_FILL_COLOR
    vector_bar.line.color.rgb = PREDICTION_VECTOR_BORDER_COLOR
    vector_bar.line.width = PREDICTION_VECTOR_BORDER_WIDTH

    # Context Vector Label (inside bar, top with equal margin)
    label_box = slide.shapes.add_textbox(
        PREDICTION_CONTENT_X + PREDICTION_VECTOR_MARGIN,
        PREDICTION_VECTOR_Y + PREDICTION_VECTOR_MARGIN,
        PREDICTION_CONTENT_WIDTH - (2 * PREDICTION_VECTOR_MARGIN), PREDICTION_VECTOR_LABEL_HEIGHT
    )
    tf = label_box.text_frame
    tf.text = "CONTEXT VECTOR"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_PREDICTION_VECTOR_LABEL
    p.font.color.rgb = FONT_COLOR_PREDICTION_VECTOR_LABEL
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_REGULAR

    # Vector Segments (10 colored rectangles, evenly distributed with equal margins)
    # Calculate available width: 12" - left margin - right margin
    available_width = PREDICTION_CONTENT_WIDTH.inches - (2 * PREDICTION_VECTOR_MARGIN.inches)
    # Calculate segment width: (available_width - 9 gaps) / 10 segments
    segment_width = (available_width - (9 * PREDICTION_SEGMENT_GAP.inches)) / PREDICTION_SEGMENT_COUNT
    # Position segments: top + top_margin + label_height + middle_margin
    segments_y = PREDICTION_VECTOR_Y + PREDICTION_VECTOR_MARGIN + PREDICTION_VECTOR_LABEL_HEIGHT + PREDICTION_VECTOR_MARGIN

    for i, color in enumerate(PREDICTION_SEGMENT_COLORS):
        segment_x = PREDICTION_CONTENT_X + PREDICTION_VECTOR_MARGIN + (i * (Inches(segment_width) + PREDICTION_SEGMENT_GAP))

        segment = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            segment_x, segments_y,
            Inches(segment_width), PREDICTION_SEGMENT_HEIGHT
        )
        segment.fill.solid()
        segment.fill.fore_color.rgb = color
        segment.line.color.rgb = color

    # Arrow down
    arrow_box = slide.shapes.add_textbox(
        PREDICTION_CONTENT_X + (PREDICTION_CONTENT_WIDTH / 2) - Inches(0.3),
        PREDICTION_ARROW_Y,
        Inches(0.6), Inches(0.5)
    )
    tf = arrow_box.text_frame
    tf.text = PREDICTION_ARROW_TEXT
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_PREDICTION_ARROW
    p.font.color.rgb = FONT_COLOR_PREDICTION_ARROW

    # Helper function to get fill color
    def get_fill_color(category):
        if category == "highest":
            return PREDICTION_FILL_HIGHEST
        elif category == "high":
            return PREDICTION_FILL_HIGH
        elif category == "medium":
            return PREDICTION_FILL_MEDIUM
        else:  # low, lowest
            return PREDICTION_FILL_LOW

    # Probability bars
    for i, pred_data in enumerate(PREDICTION_DATA):
        y_pos = PREDICTION_PROB_Y_START + (i * (PREDICTION_PROB_BAR_HEIGHT + PREDICTION_PROB_GAP))

        # Token label (left side, right-aligned text)
        label_box = slide.shapes.add_textbox(
            PREDICTION_PROB_LABEL_X, y_pos,
            PREDICTION_PROB_LABEL_WIDTH, PREDICTION_PROB_BAR_HEIGHT
        )
        tf = label_box.text_frame
        tf.text = pred_data["token"]
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = FONT_SIZE_PREDICTION_TOKEN

        # Highest token gets green color
        if pred_data["category"] == "highest":
            p.font.color.rgb = FONT_COLOR_PREDICTION_TOKEN_HIGHEST
            p.font.bold = True
        else:
            p.font.color.rgb = FONT_COLOR_PREDICTION_TOKEN

        for run in p.runs:
            run.font.name = FONT_FAMILY_PREDICTION_TOKEN

        # Probability bar background (right-aligned at fixed position)
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            PREDICTION_PROB_BAR_X, y_pos,
            PREDICTION_PROB_BAR_WIDTH, PREDICTION_PROB_BAR_HEIGHT
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = PREDICTION_PROB_BAR_FILL_COLOR
        bar_bg.line.color.rgb = PREDICTION_PROB_BAR_BORDER_COLOR
        bar_bg.line.width = PREDICTION_PROB_BAR_BORDER_WIDTH

        # Probability fill (proportional to probability, left-aligned within bar)
        fill_width = PREDICTION_PROB_BAR_WIDTH.inches * pred_data["probability"]
        if fill_width > 0.1:  # Only draw if visible
            bar_fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                PREDICTION_PROB_BAR_X, y_pos,
                Inches(fill_width), PREDICTION_PROB_BAR_HEIGHT
            )
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = get_fill_color(pred_data["category"])
            bar_fill.line.color.rgb = get_fill_color(pred_data["category"])

        # Probability value (right side of bar, always at same position)
        value_box = slide.shapes.add_textbox(
            PREDICTION_PROB_BAR_X + PREDICTION_PROB_BAR_WIDTH - Inches(0.8), y_pos,
            Inches(0.7), PREDICTION_PROB_BAR_HEIGHT
        )
        tf = value_box.text_frame
        tf.text = f"{pred_data['probability']:.2f}"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = FONT_SIZE_PREDICTION_VALUE
        p.font.color.rgb = FONT_COLOR_PREDICTION_VALUE
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_PREDICTION_VALUE

    # Thermometer Icon (left side, spans from but to no)
    # Drawn LAST so it appears in foreground (can be clicked in PowerPoint)
    thermo_pic = slide.shapes.add_picture(
        PREDICTION_THERMO_ICON,
        PREDICTION_THERMO_X, PREDICTION_THERMO_Y,
        height=PREDICTION_THERMO_HEIGHT
    )

    return prs

def create_autoregression_slide(prs, slide_num, step_data):
    """Helper function to create an autoregression slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, slide_num)

    # Title
    title_box = slide.shapes.add_textbox(
        AUTOREGRESS_TITLE_X, AUTOREGRESS_TITLE_Y,
        AUTOREGRESS_TITLE_WIDTH, AUTOREGRESS_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = step_data["title"]
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_AUTOREGRESS_TITLE
    p.font.color.rgb = FONT_COLOR_AUTOREGRESS_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_AUTOREGRESS_TITLE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        AUTOREGRESS_SUBTITLE_X, AUTOREGRESS_SUBTITLE_Y,
        AUTOREGRESS_SUBTITLE_WIDTH, AUTOREGRESS_SUBTITLE_HEIGHT
    )
    tf = subtitle_box.text_frame
    tf.text = step_data["subtitle"]
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_AUTOREGRESS_SUBTITLE
    p.font.color.rgb = FONT_COLOR_AUTOREGRESS_SUBTITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_AUTOREGRESS_SUBTITLE

    # Calculate total width of token row
    num_tokens = len(step_data["tokens"])
    tokens_width = (num_tokens * AUTOREGRESS_TOKEN_WIDTH.inches) + ((num_tokens - 1) * AUTOREGRESS_TOKEN_GAP.inches)

    # Add arrow and predicted token width if present
    if step_data["predicted"] is not None:
        tokens_width += AUTOREGRESS_LLM_ARROW_GAP.inches + AUTOREGRESS_LLM_ARROW_WIDTH.inches + AUTOREGRESS_LLM_ARROW_GAP.inches + AUTOREGRESS_TOKEN_WIDTH.inches

    # Center the entire row
    token_x_start = Inches((16 - tokens_width) / 2)

    # Token boxes row (horizontally aligned)
    for token_idx, token_text in enumerate(step_data["tokens"]):
        token_x = token_x_start + (token_idx * (AUTOREGRESS_TOKEN_WIDTH + AUTOREGRESS_TOKEN_GAP))

        # Determine if this is the new token (highlighted in green)
        is_new_token = (step_data["new_token_index"] is not None and
                      token_idx == step_data["new_token_index"])

        # Token box
        token_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            token_x, AUTOREGRESS_TOKEN_Y,
            AUTOREGRESS_TOKEN_WIDTH, AUTOREGRESS_TOKEN_HEIGHT
        )
        token_box.fill.solid()

        if is_new_token:
            # New token: green fill and border
            token_box.fill.fore_color.rgb = AUTOREGRESS_TOKEN_NEW_FILL
            token_box.line.color.rgb = AUTOREGRESS_TOKEN_NEW_BORDER
        else:
            # Existing token: blue fill and border
            token_box.fill.fore_color.rgb = AUTOREGRESS_TOKEN_FILL_COLOR
            token_box.line.color.rgb = AUTOREGRESS_TOKEN_BORDER_COLOR

        token_box.line.width = AUTOREGRESS_TOKEN_BORDER_WIDTH

        # Token text
        tf = token_box.text_frame
        tf.text = token_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_AUTOREGRESS_TOKEN

        if is_new_token:
            p.font.color.rgb = FONT_COLOR_AUTOREGRESS_TOKEN_NEW
        else:
            p.font.color.rgb = FONT_COLOR_AUTOREGRESS_TOKEN

        for run in p.runs:
            run.font.name = FONT_FAMILY_AUTOREGRESS_TOKEN

    # LLM Arrow and Predicted token (if present)
    if step_data["predicted"] is not None:
        # LLM Arrow
        llm_arrow_x = token_x_start + (num_tokens * (AUTOREGRESS_TOKEN_WIDTH + AUTOREGRESS_TOKEN_GAP)) - AUTOREGRESS_TOKEN_GAP + AUTOREGRESS_LLM_ARROW_GAP
        llm_arrow_box = slide.shapes.add_textbox(
            llm_arrow_x, AUTOREGRESS_TOKEN_Y,
            AUTOREGRESS_LLM_ARROW_WIDTH, AUTOREGRESS_TOKEN_HEIGHT
        )
        tf = llm_arrow_box.text_frame
        tf.text = AUTOREGRESS_LLM_ARROW
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_AUTOREGRESS_LLM_ARROW
        p.font.color.rgb = FONT_COLOR_AUTOREGRESS_LLM_ARROW

        # Predicted token (right of arrow, in green)
        predicted_x = llm_arrow_x + AUTOREGRESS_LLM_ARROW_WIDTH + AUTOREGRESS_LLM_ARROW_GAP
        predicted_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            predicted_x, AUTOREGRESS_TOKEN_Y,
            AUTOREGRESS_TOKEN_WIDTH, AUTOREGRESS_TOKEN_HEIGHT
        )
        predicted_box.fill.solid()
        predicted_box.fill.fore_color.rgb = AUTOREGRESS_TOKEN_NEW_FILL
        predicted_box.line.color.rgb = AUTOREGRESS_TOKEN_NEW_BORDER
        predicted_box.line.width = AUTOREGRESS_TOKEN_BORDER_WIDTH

        tf = predicted_box.text_frame
        tf.text = step_data["predicted"]
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_AUTOREGRESS_TOKEN
        p.font.color.rgb = FONT_COLOR_AUTOREGRESS_TOKEN_NEW
        for run in p.runs:
            run.font.name = FONT_FAMILY_AUTOREGRESS_TOKEN
    # Final slide: no completion message (user requested clean final slide)

    return prs

def create_slide_12(prs):
    """Slide 12: Autoregression - Step 1"""
    return create_autoregression_slide(prs, 12, AUTOREGRESS_STEP_1)

def create_slide_13(prs):
    """Slide 13: Autoregression - Step 2"""
    return create_autoregression_slide(prs, 13, AUTOREGRESS_STEP_2)

def create_slide_14(prs):
    """Slide 14: Autoregression - Step 3"""
    return create_autoregression_slide(prs, 14, AUTOREGRESS_STEP_3)

def create_slide_15(prs):
    """Slide 15: Autoregression - Final"""
    return create_autoregression_slide(prs, 15, AUTOREGRESS_STEP_FINAL)

def create_slide_16(prs):
    """Slide 16: ON PREMISE MATTERS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 16, 25)

    # The three keywords - using KEYWORD_THEME_SOLUTION
    keywords = [
        {"text": "ON", "color": KEYWORD_THEME_SOLUTION[0]},
        {"text": "PREMISE", "color": KEYWORD_THEME_SOLUTION[1]},
        {"text": "MATTERS", "color": KEYWORD_THEME_SOLUTION[2]}
    ]

    for i, keyword in enumerate(keywords):
        y_pos = KEYWORD_Y_START + (i * KEYWORD_Y_GAP)

        keyword_box = slide.shapes.add_textbox(
            KEYWORD_BOX_X, Inches(y_pos),
            KEYWORD_BOX_WIDTH, KEYWORD_BOX_HEIGHT
        )
        tf = keyword_box.text_frame
        tf.text = keyword["text"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_KEYWORD
        p.font.bold = FONT_BOLD_KEYWORD
        p.font.color.rgb = keyword["color"]

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_KEYWORD  # Inter ExtraLight (font-weight: 200)
            run.font.character_spacing = FONT_LETTER_SPACING_KEYWORD

    return prs

def create_slide_17(prs):
    """Slide 17: The Fundamental Security Conflict"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 17, 25)

    # Title
    title_box = slide.shapes.add_textbox(
        SECURITY_TITLE_X, SECURITY_TITLE_Y,
        SECURITY_TITLE_WIDTH, SECURITY_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "The Fundamental Security Conflict"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_SECURITY_TITLE
    p.font.color.rgb = FONT_COLOR_SECURITY_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_SECURITY_TITLE

    # === LEFT CARD: Cloud Providers (Red) ===
    # Card container with red border
    cloud_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SECURITY_CARD_LEFT_X, SECURITY_CARD_Y,
        SECURITY_CARD_WIDTH, SECURITY_CARD_HEIGHT
    )
    cloud_card.fill.solid()
    cloud_card.fill.fore_color.rgb = SECURITY_CARD_FILL_COLOR
    cloud_card.line.color.rgb = COLOR_SECURITY_CLOUD
    cloud_card.line.width = SECURITY_CARD_BORDER_WIDTH
    cloud_card.adjustments[0] = 0.05  # Rounded corners

    # Card title
    cloud_title_box = slide.shapes.add_textbox(
        SECURITY_CARD_LEFT_X, SECURITY_CARD_Y + Inches(SECURITY_COL_TITLE_Y_OFFSET),
        SECURITY_CARD_WIDTH, SECURITY_COL_TITLE_HEIGHT
    )
    tf = cloud_title_box.text_frame
    tf.text = "Cloud Providers"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_SECURITY_COL_TITLE
    p.font.color.rgb = COLOR_SECURITY_CLOUD
    p.font.bold = True
    for run in p.runs:
        run.font.name = FONT_FAMILY_SECURITY_COL_TITLE

    # Cloud icon
    cloud_icon_x = SECURITY_CARD_LEFT_X + (SECURITY_CARD_WIDTH - SECURITY_ICON_WIDTH) / 2
    cloud_icon_y = SECURITY_CARD_Y + Inches(SECURITY_ICON_Y_OFFSET)
    cloud_icon = slide.shapes.add_picture(
        SECURITY_CLOUD_ICON,
        cloud_icon_x,
        cloud_icon_y,
        width=SECURITY_ICON_WIDTH
    )

    # Cloud steps
    for i, step_text in enumerate(SECURITY_CLOUD_STEPS):
        step_y = SECURITY_CARD_Y + Inches(SECURITY_STEP_Y_START_OFFSET + (i * (SECURITY_STEP_HEIGHT.inches + SECURITY_STEP_GAP.inches)))

        # Step number in colored circle
        num_box = slide.shapes.add_textbox(
            SECURITY_CARD_LEFT_X + Inches(SECURITY_STEP_X_OFFSET), step_y,
            Inches(0.5), SECURITY_STEP_HEIGHT
        )
        tf = num_box.text_frame
        tf.text = str(i + 1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_SECURITY_STEP_NUMBER
        p.font.color.rgb = COLOR_SECURITY_CLOUD
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_SECURITY_STEP

        # Step text
        text_box = slide.shapes.add_textbox(
            SECURITY_CARD_LEFT_X + Inches(SECURITY_STEP_X_OFFSET + 0.7), step_y,
            SECURITY_CARD_WIDTH - Inches(SECURITY_STEP_X_OFFSET + 0.9), SECURITY_STEP_HEIGHT
        )
        tf = text_box.text_frame
        tf.text = step_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = FONT_SIZE_SECURITY_STEP_TEXT
        p.font.color.rgb = COLOR_TEXT_WHITE
        for run in p.runs:
            run.font.name = FONT_FAMILY_SECURITY_STEP

    # === RIGHT CARD: Brain-Bridges (Green) ===
    # Card container with green border
    local_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SECURITY_CARD_RIGHT_X, SECURITY_CARD_Y,
        SECURITY_CARD_WIDTH, SECURITY_CARD_HEIGHT
    )
    local_card.fill.solid()
    local_card.fill.fore_color.rgb = SECURITY_CARD_FILL_COLOR
    local_card.line.color.rgb = COLOR_SECURITY_LOCAL
    local_card.line.width = SECURITY_CARD_BORDER_WIDTH
    local_card.adjustments[0] = 0.05  # Rounded corners

    # Card title
    local_title_box = slide.shapes.add_textbox(
        SECURITY_CARD_RIGHT_X, SECURITY_CARD_Y + Inches(SECURITY_COL_TITLE_Y_OFFSET),
        SECURITY_CARD_WIDTH, SECURITY_COL_TITLE_HEIGHT
    )
    tf = local_title_box.text_frame
    tf.text = "Brain-Bridges"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_SECURITY_COL_TITLE
    p.font.color.rgb = COLOR_SECURITY_LOCAL
    p.font.bold = True
    for run in p.runs:
        run.font.name = FONT_FAMILY_SECURITY_COL_TITLE

    # Local icon
    local_icon_x = SECURITY_CARD_RIGHT_X + (SECURITY_CARD_WIDTH - SECURITY_ICON_WIDTH) / 2
    local_icon_y = SECURITY_CARD_Y + Inches(SECURITY_ICON_Y_OFFSET)
    local_icon = slide.shapes.add_picture(
        SECURITY_LOCAL_ICON,
        local_icon_x,
        local_icon_y,
        width=SECURITY_ICON_WIDTH
    )

    # Local steps
    for i, step_text in enumerate(SECURITY_LOCAL_STEPS):
        step_y = SECURITY_CARD_Y + Inches(SECURITY_STEP_Y_START_OFFSET + (i * (SECURITY_STEP_HEIGHT.inches + SECURITY_STEP_GAP.inches)))

        # Step number in colored circle
        num_box = slide.shapes.add_textbox(
            SECURITY_CARD_RIGHT_X + Inches(SECURITY_STEP_X_OFFSET), step_y,
            Inches(0.5), SECURITY_STEP_HEIGHT
        )
        tf = num_box.text_frame
        tf.text = str(i + 1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_SECURITY_STEP_NUMBER
        p.font.color.rgb = COLOR_SECURITY_LOCAL
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_SECURITY_STEP

        # Step text
        text_box = slide.shapes.add_textbox(
            SECURITY_CARD_RIGHT_X + Inches(SECURITY_STEP_X_OFFSET + 0.7), step_y,
            SECURITY_CARD_WIDTH - Inches(SECURITY_STEP_X_OFFSET + 0.9), SECURITY_STEP_HEIGHT
        )
        tf = text_box.text_frame
        tf.text = step_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = FONT_SIZE_SECURITY_STEP_TEXT
        p.font.color.rgb = COLOR_TEXT_WHITE
        for run in p.runs:
            run.font.name = FONT_FAMILY_SECURITY_STEP

    return prs

def create_slide_18(prs):
    """Slide 18: The Encryption Dilemma - Redesigned Layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 18, 25)

    # Title
    title_box = slide.shapes.add_textbox(
        ENCRYPTION_TITLE_X, ENCRYPTION_TITLE_Y,
        ENCRYPTION_TITLE_WIDTH, ENCRYPTION_TITLE_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = "The Encryption Dilemma"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_TITLE
    p.font.color.rgb = FONT_COLOR_ENCRYPTION_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_TITLE

    # Subtitle (red mono font like Slide 2)
    subtitle_box = slide.shapes.add_textbox(
        ENCRYPTION_TITLE_X, ENCRYPTION_SUBTITLE_Y,
        ENCRYPTION_TITLE_WIDTH, ENCRYPTION_TITLE_HEIGHT
    )
    tf = subtitle_box.text_frame
    tf.text = "Data must be decrypted for inference processing"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_SUBTITLE
    p.font.color.rgb = FONT_COLOR_CONTENT_SUBTITLE_ALERT  # Red like Slide 2
    for run in p.runs:
        run.font.name = FONT_FAMILY_SUBTITLE  # Menlo mono font

    # === STEP 3: Remote Cloud Server (Top Center) ===
    # Card
    inference_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        ENCRYPTION_TOP_X, ENCRYPTION_TOP_Y,
        ENCRYPTION_TOP_CARD_WIDTH, ENCRYPTION_TOP_CARD_HEIGHT
    )
    inference_card.fill.solid()
    inference_card.fill.fore_color.rgb = ENCRYPTION_CARD_FILL_COLOR
    inference_card.line.color.rgb = COLOR_ENCRYPTION_INFERENCE
    inference_card.line.width = ENCRYPTION_CARD_BORDER_WIDTH
    inference_card.adjustments[0] = 0.05

    # Icon (drawn first so title text appears on top)
    icon_x = ENCRYPTION_TOP_X + Inches(ENCRYPTION_CLOUD_ICON_X_OFFSET)
    icon_y = ENCRYPTION_TOP_Y + Inches(ENCRYPTION_ICON_Y_OFFSET)
    cloud_icon = slide.shapes.add_picture(
        ENCRYPTION_CLOUD_ICON,
        icon_x, icon_y,
        width=ENCRYPTION_CLOUD_ICON_WIDTH
    )

    # Title (drawn after icon, appears on top in vertical layout, centered)
    stage_title_box = slide.shapes.add_textbox(
        ENCRYPTION_TOP_X + Inches(ENCRYPTION_STAGE_TITLE_X_OFFSET),
        ENCRYPTION_TOP_Y + Inches(ENCRYPTION_STAGE_TITLE_Y_OFFSET),
        ENCRYPTION_TOP_CARD_WIDTH - Inches(0.6), Inches(0.4)
    )
    tf = stage_title_box.text_frame
    tf.text = "Remote Cloud Server"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_STAGE_TITLE
    p.font.color.rgb = COLOR_ENCRYPTION_INFERENCE
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_STAGE_TITLE

    # Step number (blue circle, drawn last to appear on top)
    step3_num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        ENCRYPTION_TOP_X + Inches(ENCRYPTION_STEP_NUMBER_X_OFFSET_INSIDE),
        ENCRYPTION_TOP_Y + Inches(ENCRYPTION_STEP_NUMBER_Y_OFFSET_INSIDE),
        ENCRYPTION_STEP_NUMBER_SIZE, ENCRYPTION_STEP_NUMBER_SIZE
    )
    step3_num_circle.fill.solid()
    step3_num_circle.fill.fore_color.rgb = COLOR_ENCRYPTION_STEP_NUMBER_BG
    step3_num_circle.line.width = Pt(0)
    # Text directly in the OVAL shape
    tf = step3_num_circle.text_frame
    tf.text = "3"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_STEP_NUMBER
    p.font.color.rgb = COLOR_ENCRYPTION_STEP_NUMBER_TEXT
    p.font.bold = True
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_STAGE_TITLE

    # Processing text
    processing_box = slide.shapes.add_textbox(
        ENCRYPTION_TOP_X + Inches(0.5),
        ENCRYPTION_TOP_Y + Inches(ENCRYPTION_PROCESSING_Y_OFFSET),
        ENCRYPTION_TOP_CARD_WIDTH - Inches(1.0), Inches(0.6)
    )
    tf = processing_box.text_frame
    tf.text = "Inferencing"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_PROCESSING
    p.font.color.rgb = COLOR_ENCRYPTION_INFERENCE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_REGULAR

    # Description - REMOVED (user request)
    # desc_box = slide.shapes.add_textbox(
    #     ENCRYPTION_TOP_X + Inches(0.3),
    #     ENCRYPTION_TOP_Y + Inches(ENCRYPTION_DESC_Y_OFFSET) - Inches(0.5),
    #     ENCRYPTION_TOP_CARD_WIDTH - Inches(0.6), Inches(0.4)
    # )
    # tf = desc_box.text_frame
    # tf.text = "LLM processes plain text"
    # p = tf.paragraphs[0]
    # p.alignment = PP_ALIGN.CENTER
    # p.font.size = FONT_SIZE_ENCRYPTION_DESC
    # p.font.color.rgb = COLOR_TEXT_GRAY
    # for run in p.runs:
    #     run.font.name = FONT_FAMILY_INTER_REGULAR

    # === STEP 1: Encrypted Data (Bottom Left) ===
    # Card
    encrypted_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        ENCRYPTION_LEFT_X, ENCRYPTION_BOTTOM_Y,
        ENCRYPTION_CARD_WIDTH, ENCRYPTION_CARD_HEIGHT
    )
    encrypted_card.fill.solid()
    encrypted_card.fill.fore_color.rgb = ENCRYPTION_CARD_FILL_COLOR
    encrypted_card.line.color.rgb = COLOR_ENCRYPTION_ENCRYPTED
    encrypted_card.line.width = ENCRYPTION_CARD_BORDER_WIDTH
    encrypted_card.adjustments[0] = 0.05

    # Icon (drawn first so title text appears on top)
    icon_x = ENCRYPTION_LEFT_X + Inches(ENCRYPTION_LOCK_ICON_X_OFFSET)
    icon_y = ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_ICON_Y_OFFSET)
    lock_icon = slide.shapes.add_picture(
        ENCRYPTION_LOCK_ICON,
        icon_x, icon_y,
        width=ENCRYPTION_LOCK_ICON_WIDTH
    )

    # Title (drawn after icon, appears on top in vertical layout, centered)
    stage_title_box = slide.shapes.add_textbox(
        ENCRYPTION_LEFT_X + Inches(ENCRYPTION_STAGE_TITLE_X_OFFSET),
        ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_STAGE_TITLE_Y_OFFSET),
        ENCRYPTION_CARD_WIDTH - Inches(0.6), Inches(0.4)
    )
    tf = stage_title_box.text_frame
    tf.text = "Encrypted Data"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_STAGE_TITLE
    p.font.color.rgb = COLOR_ENCRYPTION_ENCRYPTED
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_STAGE_TITLE

    # Step number (blue circle, drawn last to appear on top)
    step1_num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        ENCRYPTION_LEFT_X + Inches(ENCRYPTION_STEP_NUMBER_X_OFFSET_INSIDE),
        ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_STEP_NUMBER_Y_OFFSET_INSIDE),
        ENCRYPTION_STEP_NUMBER_SIZE, ENCRYPTION_STEP_NUMBER_SIZE
    )
    step1_num_circle.fill.solid()
    step1_num_circle.fill.fore_color.rgb = COLOR_ENCRYPTION_STEP_NUMBER_BG
    step1_num_circle.line.width = Pt(0)
    # Text directly in the OVAL shape
    tf = step1_num_circle.text_frame
    tf.text = "1"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_STEP_NUMBER
    p.font.color.rgb = COLOR_ENCRYPTION_STEP_NUMBER_TEXT
    p.font.bold = True
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_STAGE_TITLE

    # Data blocks (encrypted) - styled with borders
    for i, block_text in enumerate(ENCRYPTION_ENCRYPTED_BLOCKS):
        row = i // 3
        col = i % 3
        block_x = ENCRYPTION_LEFT_X + Inches(ENCRYPTION_DATA_START_X_OFFSET) + (col * (ENCRYPTION_DATA_BLOCK_WIDTH + ENCRYPTION_DATA_BLOCK_GAP_X))
        block_y = ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_DATA_Y_OFFSET) + (row * (ENCRYPTION_DATA_BLOCK_HEIGHT + ENCRYPTION_DATA_BLOCK_GAP_Y))

        # Box with border and text inside
        block_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            block_x, block_y,
            ENCRYPTION_DATA_BLOCK_WIDTH, ENCRYPTION_DATA_BLOCK_HEIGHT
        )
        block_shape.fill.solid()
        block_shape.fill.fore_color.rgb = ENCRYPTION_CARD_FILL_COLOR
        block_shape.line.color.rgb = COLOR_ENCRYPTION_ENCRYPTED
        block_shape.line.width = ENCRYPTION_DATA_BLOCK_BORDER_WIDTH
        block_shape.adjustments[0] = 0.1

        # Text directly in the shape
        tf = block_shape.text_frame
        tf.text = block_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_ENCRYPTION_DATA
        p.font.color.rgb = COLOR_ENCRYPTION_ENCRYPTED
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

    # Description - REMOVED (user request)
    # desc_box = slide.shapes.add_textbox(
    #     ENCRYPTION_LEFT_X + Inches(0.3),
    #     ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_DESC_Y_OFFSET),
    #     ENCRYPTION_CARD_WIDTH - Inches(0.6), Inches(0.4)
    # )
    # tf = desc_box.text_frame
    # tf.text = "Secure network transmission"
    # p = tf.paragraphs[0]
    # p.alignment = PP_ALIGN.CENTER
    # p.font.size = FONT_SIZE_ENCRYPTION_DESC
    # p.font.color.rgb = COLOR_TEXT_GRAY
    # for run in p.runs:
    #     run.font.name = FONT_FAMILY_INTER_REGULAR

    # === STEP 2: Decrypted Data (Bottom Right) ===
    # Card
    decrypted_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        ENCRYPTION_RIGHT_X, ENCRYPTION_BOTTOM_Y,
        ENCRYPTION_CARD_WIDTH, ENCRYPTION_CARD_HEIGHT
    )
    decrypted_card.fill.solid()
    decrypted_card.fill.fore_color.rgb = ENCRYPTION_CARD_FILL_COLOR
    decrypted_card.line.color.rgb = COLOR_ENCRYPTION_DECRYPTED
    decrypted_card.line.width = ENCRYPTION_CARD_BORDER_WIDTH
    decrypted_card.adjustments[0] = 0.05

    # Icon (drawn first so title text appears on top)
    icon_x = ENCRYPTION_RIGHT_X + Inches(ENCRYPTION_LOCK_ICON_X_OFFSET)
    icon_y = ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_ICON_Y_OFFSET)
    unlock_icon = slide.shapes.add_picture(
        ENCRYPTION_UNLOCK_ICON,
        icon_x, icon_y,
        width=ENCRYPTION_LOCK_ICON_WIDTH
    )

    # Title (drawn after icon, appears on top in vertical layout, centered)
    stage_title_box = slide.shapes.add_textbox(
        ENCRYPTION_RIGHT_X + Inches(ENCRYPTION_STAGE_TITLE_X_OFFSET),
        ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_STAGE_TITLE_Y_OFFSET),
        ENCRYPTION_CARD_WIDTH - Inches(0.6), Inches(0.4)
    )
    tf = stage_title_box.text_frame
    tf.text = "Decrypted Data"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_STAGE_TITLE
    p.font.color.rgb = COLOR_ENCRYPTION_DECRYPTED
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_STAGE_TITLE

    # Step number (blue circle, drawn last to appear on top)
    step2_num_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        ENCRYPTION_RIGHT_X + Inches(ENCRYPTION_STEP_NUMBER_X_OFFSET_INSIDE),
        ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_STEP_NUMBER_Y_OFFSET_INSIDE),
        ENCRYPTION_STEP_NUMBER_SIZE, ENCRYPTION_STEP_NUMBER_SIZE
    )
    step2_num_circle.fill.solid()
    step2_num_circle.fill.fore_color.rgb = COLOR_ENCRYPTION_STEP_NUMBER_BG
    step2_num_circle.line.width = Pt(0)
    # Text directly in the OVAL shape
    tf = step2_num_circle.text_frame
    tf.text = "2"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_ENCRYPTION_STEP_NUMBER
    p.font.color.rgb = COLOR_ENCRYPTION_STEP_NUMBER_TEXT
    p.font.bold = True
    for run in p.runs:
        run.font.name = FONT_FAMILY_ENCRYPTION_STAGE_TITLE

    # Data blocks (decrypted) - styled with borders
    for i, block_text in enumerate(ENCRYPTION_DECRYPTED_BLOCKS):
        row = i // 3
        col = i % 3
        block_x = ENCRYPTION_RIGHT_X + Inches(ENCRYPTION_DATA_START_X_OFFSET) + (col * (ENCRYPTION_DATA_BLOCK_WIDTH + ENCRYPTION_DATA_BLOCK_GAP_X))
        block_y = ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_DATA_Y_OFFSET) + (row * (ENCRYPTION_DATA_BLOCK_HEIGHT + ENCRYPTION_DATA_BLOCK_GAP_Y))

        # Box with border and text inside
        block_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            block_x, block_y,
            ENCRYPTION_DATA_BLOCK_WIDTH, ENCRYPTION_DATA_BLOCK_HEIGHT
        )
        block_shape.fill.solid()
        block_shape.fill.fore_color.rgb = ENCRYPTION_CARD_FILL_COLOR
        block_shape.line.color.rgb = COLOR_ENCRYPTION_DECRYPTED
        block_shape.line.width = ENCRYPTION_DATA_BLOCK_BORDER_WIDTH
        block_shape.adjustments[0] = 0.1

        # Text directly in the shape
        tf = block_shape.text_frame
        tf.text = block_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_ENCRYPTION_DATA
        p.font.color.rgb = COLOR_ENCRYPTION_DECRYPTED
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

    # Description - REMOVED (user request)
    # desc_box = slide.shapes.add_textbox(
    #     ENCRYPTION_RIGHT_X + Inches(0.3),
    #     ENCRYPTION_BOTTOM_Y + Inches(ENCRYPTION_DESC_Y_OFFSET),
    #     ENCRYPTION_CARD_WIDTH - Inches(0.6), Inches(0.4)
    # )
    # tf = desc_box.text_frame
    # tf.text = "Plain text for AI processing"
    # p = tf.paragraphs[0]
    # p.alignment = PP_ALIGN.CENTER
    # p.font.size = FONT_SIZE_ENCRYPTION_DESC
    # p.font.color.rgb = COLOR_TEXT_GRAY
    # for run in p.runs:
    #     run.font.name = FONT_FAMILY_INTER_REGULAR

    return prs
def create_placeholder_slide(prs, slide_num):
    """Creates a placeholder slide for later editing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, slide_num)

    # Placeholder title
    title_box = slide.shapes.add_textbox(
        PLACEHOLDER_X, PLACEHOLDER_Y,
        PLACEHOLDER_WIDTH, PLACEHOLDER_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = f"Slide {slide_num}\n(To be designed)"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_PLACEHOLDER
    p.font.color.rgb = FONT_COLOR_PLACEHOLDER

    return prs

def create_slide_19(prs):
    """Slide 19: Chat API Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 19, 25)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), CHAT_API_TITLE_Y,
        Inches(14), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.text = "Chat API Architecture"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CHAT_API_TITLE
    p.font.color.rgb = FONT_COLOR_CHAT_API_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_EXTRALIGHT

    # Card data: (role_name, badge_text, note_text, content_text, color)
    cards = [
        ("SYSTEM", "SYSTEM", "Hidden instructions that guide AI behavior", CHAT_API_SYSTEM_TEXT, COLOR_CHAT_API_SYSTEM),
        ("USER", "USER", "Actual request from the user", CHAT_API_USER_TEXT, COLOR_CHAT_API_USER),
        ("ASSISTANT", "ASSISTANT", "AI response following system instructions", CHAT_API_ASSISTANT_TEXT, COLOR_CHAT_API_ASSISTANT)
    ]

    # Create each card
    for i, (role_name, badge_text, note_text, content_text, color) in enumerate(cards):
        card_y = CHAT_API_CARD_START_Y + (i * (CHAT_API_CARD_HEIGHT + CHAT_API_CARD_GAP))

        # Card frame - NO BACKGROUND, only border (minimalist)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CHAT_API_CARD_X, card_y,
            CHAT_API_CARD_WIDTH, CHAT_API_CARD_HEIGHT
        )
        # No background fill - transparent
        card.fill.background()
        # Colored border only
        card.line.color.rgb = color
        card.line.width = CHAT_API_CARD_BORDER_WIDTH
        card.adjustments[0] = CHAT_API_CARD_BORDER_RADIUS
        # No shadow
        card.shadow.inherit = False

        # Role Badge - NO BACKGROUND, only border (minimalist)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CHAT_API_CARD_X + CHAT_API_METADATA_X_OFFSET,
            card_y + CHAT_API_BADGE_Y_OFFSET,
            CHAT_API_BADGE_WIDTH, CHAT_API_BADGE_HEIGHT
        )
        # No background fill - transparent
        badge.fill.background()
        # Colored border only
        badge.line.color.rgb = color
        badge.line.width = CHAT_API_BADGE_BORDER_WIDTH
        badge.adjustments[0] = 0.5  # Very rounded for pill shape
        # No shadow
        badge.shadow.inherit = False

        # Badge text (full color)
        tf = badge.text_frame
        tf.text = badge_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_CHAT_API_BADGE
        p.font.color.rgb = color
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

        # Message note (description under badge)
        note_box = slide.shapes.add_textbox(
            CHAT_API_CARD_X + CHAT_API_METADATA_X_OFFSET,
            card_y + CHAT_API_NOTE_Y_OFFSET,
            CHAT_API_METADATA_WIDTH, Inches(0.4)
        )
        tf = note_box.text_frame
        tf.text = note_text
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = FONT_SIZE_CHAT_API_NOTE
        p.font.color.rgb = FONT_COLOR_CHAT_API_NOTE
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

        # Vertical divider line (using a thin rectangle shape)
        divider_height = CHAT_API_CARD_HEIGHT - (2 * CHAT_API_DIVIDER_HEIGHT_OFFSET)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            CHAT_API_CARD_X + CHAT_API_DIVIDER_X_OFFSET,
            card_y + CHAT_API_DIVIDER_HEIGHT_OFFSET,
            Pt(2),  # Width
            divider_height
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(255, 255, 255)
        divider.fill.transparency = 0.7  # Semi-transparent
        divider.line.fill.background()  # No border

        # Message content (main text on right)
        content_box = slide.shapes.add_textbox(
            CHAT_API_CARD_X + CHAT_API_CONTENT_X_OFFSET,
            card_y + Inches(0.15),  # Centered vertically
            CHAT_API_CONTENT_WIDTH,
            CHAT_API_CARD_HEIGHT - Inches(0.3)
        )
        tf = content_box.text_frame
        tf.text = content_text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = FONT_SIZE_CHAT_API_CONTENT
        p.font.color.rgb = FONT_COLOR_CHAT_API_CONTENT
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

    return prs

def create_slide_20(prs):
    """Slide 20: Chat API Architecture (copy of slide 19)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 20, 25)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), CHAT_API_TITLE_Y,
        Inches(14), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.text = "Chat API Architecture"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CHAT_API_TITLE
    p.font.color.rgb = FONT_COLOR_CHAT_API_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_EXTRALIGHT

    # Card data: (role_name, badge_text, note_text, content_text, color)
    cards = [
        ("SYSTEM", "SYSTEM", "Hidden instructions that guide AI behavior", CHAT_API_SYSTEM_TEXT, COLOR_CHAT_API_SYSTEM),
        ("USER", "USER", "Actual request from the user", CHAT_API_USER_TEXT, COLOR_CHAT_API_USER),
        ("ASSISTANT", "ASSISTANT", "AI response following system instructions", CHAT_API_ASSISTANT_TEXT, COLOR_CHAT_API_ASSISTANT)
    ]

    # Create each card
    for i, (role_name, badge_text, note_text, content_text, color) in enumerate(cards):
        card_y = CHAT_API_CARD_START_Y + (i * (CHAT_API_CARD_HEIGHT + CHAT_API_CARD_GAP))

        # Card frame - NO BACKGROUND, only border (minimalist)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CHAT_API_CARD_X, card_y,
            CHAT_API_CARD_WIDTH, CHAT_API_CARD_HEIGHT
        )
        card.fill.background()
        card.line.color.rgb = color
        card.line.width = CHAT_API_CARD_BORDER_WIDTH
        card.adjustments[0] = CHAT_API_CARD_BORDER_RADIUS
        card.shadow.inherit = False

        # Role Badge - NO BACKGROUND, only border (minimalist)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            CHAT_API_CARD_X + CHAT_API_METADATA_X_OFFSET,
            card_y + CHAT_API_BADGE_Y_OFFSET,
            CHAT_API_BADGE_WIDTH, CHAT_API_BADGE_HEIGHT
        )
        badge.fill.background()
        badge.line.color.rgb = color
        badge.line.width = CHAT_API_BADGE_BORDER_WIDTH
        badge.adjustments[0] = 0.5
        badge.shadow.inherit = False

        # Badge text
        tf = badge.text_frame
        tf.text = badge_text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_CHAT_API_BADGE
        p.font.color.rgb = color
        p.font.bold = True
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

        # Message note
        note_box = slide.shapes.add_textbox(
            CHAT_API_CARD_X + CHAT_API_METADATA_X_OFFSET,
            card_y + CHAT_API_NOTE_Y_OFFSET,
            CHAT_API_METADATA_WIDTH, Inches(0.4)
        )
        tf = note_box.text_frame
        tf.text = note_text
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = FONT_SIZE_CHAT_API_NOTE
        p.font.color.rgb = FONT_COLOR_CHAT_API_NOTE
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

        # Vertical divider
        divider_height = CHAT_API_CARD_HEIGHT - (2 * CHAT_API_DIVIDER_HEIGHT_OFFSET)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            CHAT_API_CARD_X + CHAT_API_DIVIDER_X_OFFSET,
            card_y + CHAT_API_DIVIDER_HEIGHT_OFFSET,
            Pt(2),
            divider_height
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(255, 255, 255)
        divider.fill.transparency = 0.7
        divider.line.fill.background()

        # Message content
        content_box = slide.shapes.add_textbox(
            CHAT_API_CARD_X + CHAT_API_CONTENT_X_OFFSET,
            card_y + Inches(0.15),
            CHAT_API_CONTENT_WIDTH,
            CHAT_API_CARD_HEIGHT - Inches(0.3)
        )
        tf = content_box.text_frame
        tf.text = content_text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = FONT_SIZE_CHAT_API_CONTENT
        p.font.color.rgb = FONT_COLOR_CHAT_API_CONTENT
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_REGULAR

    return prs

def create_slide_21(prs):
    """Slide 21: RETRIEVAL AUGMENTED GENERATION"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 21, 25)

    # The three keywords - using KEYWORD_THEME_TECH (like slides 4, 7, 16)
    keywords = [
        {"text": "RETRIEVAL", "color": KEYWORD_THEME_TECH[0]},
        {"text": "AUGMENTED", "color": KEYWORD_THEME_TECH[1]},
        {"text": "GENERATION", "color": KEYWORD_THEME_TECH[2]}
    ]

    for i, keyword in enumerate(keywords):
        y_pos = KEYWORD_Y_START + (i * KEYWORD_Y_GAP)

        keyword_box = slide.shapes.add_textbox(
            KEYWORD_BOX_X, Inches(y_pos),
            KEYWORD_BOX_WIDTH, KEYWORD_BOX_HEIGHT
        )
        tf = keyword_box.text_frame
        tf.text = keyword["text"]
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_KEYWORD
        p.font.bold = FONT_BOLD_KEYWORD
        p.font.color.rgb = keyword["color"]

        # CRITICAL: Font name must be set at RUN level!
        for run in p.runs:
            run.font.name = FONT_FAMILY_KEYWORD  # Inter ExtraLight (font-weight: 200)
            run.font.character_spacing = FONT_LETTER_SPACING_KEYWORD

    return prs

def create_slide_22(prs):
    """Slide 22: Document Processing (RAG)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, 22, 25)

    # Title - Fixed header
    title_box = slide.shapes.add_textbox(
        CONTENT_HEADER_X, CONTENT_HEADER_Y,
        CONTENT_HEADER_WIDTH, CONTENT_HEADER_HEIGHT
    )
    tf = title_box.text_frame
    tf.text = DOC_PROC_TITLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_CONTENT_TITLE
    p.font.bold = FONT_BOLD_CONTENT_TITLE
    p.font.color.rgb = FONT_COLOR_CONTENT_TITLE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_SEMIBOLD

    # === LEFT: PDF Documents ===
    # PDF1 (red) - bottom layer
    pdf1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        DOC_PROC_PDF_X, DOC_PROC_PDF_Y,
        DOC_PROC_PDF_WIDTH, DOC_PROC_PDF_HEIGHT
    )
    pdf1.fill.solid()
    pdf1.fill.fore_color.rgb = COLOR_BACKGROUND_LIGHT
    pdf1.line.color.rgb = RGBColor(64, 64, 64)
    pdf1.line.width = DOC_PROC_PDF_BORDER_WIDTH
    pdf1.rotation = -5  # Slight rotation

    # PDF1 Header
    pdf1_header = slide.shapes.add_textbox(
        DOC_PROC_PDF_X + Inches(0.15), DOC_PROC_PDF_Y + Inches(0.15),
        DOC_PROC_PDF_WIDTH - Inches(0.3), Inches(0.4)
    )
    tf = pdf1_header.text_frame
    tf.text = "📄 PDF1"
    p = tf.paragraphs[0]
    p.font.size = FONT_SIZE_DOC_PROC_PDF_HEADER
    p.font.bold = True
    p.font.color.rgb = COLOR_PDF1
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_SEMIBOLD

    # PDF2 (cyan) - top layer, offset
    pdf2_x = DOC_PROC_PDF_X + Inches(0.4)
    pdf2_y = DOC_PROC_PDF_Y + Inches(0.3)
    pdf2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        pdf2_x, pdf2_y,
        DOC_PROC_PDF_WIDTH, DOC_PROC_PDF_HEIGHT
    )
    pdf2.fill.solid()
    pdf2.fill.fore_color.rgb = COLOR_BACKGROUND_LIGHT
    pdf2.line.color.rgb = RGBColor(64, 64, 64)
    pdf2.line.width = DOC_PROC_PDF_BORDER_WIDTH
    pdf2.rotation = 3  # Slight rotation opposite direction

    # PDF2 Header
    pdf2_header = slide.shapes.add_textbox(
        pdf2_x + Inches(0.15), pdf2_y + Inches(0.15),
        DOC_PROC_PDF_WIDTH - Inches(0.3), Inches(0.4)
    )
    tf = pdf2_header.text_frame
    tf.text = "📄 PDF2"
    p = tf.paragraphs[0]
    p.font.size = FONT_SIZE_DOC_PROC_PDF_HEADER
    p.font.bold = True
    p.font.color.rgb = COLOR_PDF2
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_SEMIBOLD

    # === CENTER: Arrow ===
    arrow_box = slide.shapes.add_textbox(
        DOC_PROC_ARROW_X, DOC_PROC_ARROW_Y,
        Inches(1), Inches(1)
    )
    tf = arrow_box.text_frame
    tf.text = "→"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = DOC_PROC_ARROW_SIZE
    p.font.color.rgb = COLOR_ACCENT_BLUE

    # === RIGHT: Vector Matrix ===
    chunks = [
        ("Chunk 1", ["0.23", "-0.15", "0.87", "-0.42", "0.66", "..."]),
        ("Chunk 2", ["0.45", "-0.67", "0.12", "0.89", "-0.34", "..."]),
        ("Chunk 3", ["-0.56", "0.78", "-0.23", "0.41", "0.92", "..."]),
        ("Chunk n", ["0.33", "-0.42", "0.71", "-0.18", "0.64", "..."])
    ]

    current_y = DOC_PROC_VECTOR_Y
    for chunk_label, values in chunks:
        # Label
        label_box = slide.shapes.add_textbox(
            DOC_PROC_VECTOR_X, current_y,
            DOC_PROC_VECTOR_LABEL_WIDTH, DOC_PROC_VECTOR_CELL_HEIGHT
        )
        tf = label_box.text_frame
        tf.text = chunk_label
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = FONT_SIZE_DOC_PROC_LABEL
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_BLUE
        for run in p.runs:
            run.font.name = FONT_FAMILY_INTER_SEMIBOLD

        # Vector cells
        cell_x = DOC_PROC_VECTOR_X + DOC_PROC_VECTOR_LABEL_WIDTH + Inches(0.2)
        for i, value in enumerate(values):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                cell_x + (i * (DOC_PROC_VECTOR_CELL_WIDTH + Inches(0.08))), current_y,
                DOC_PROC_VECTOR_CELL_WIDTH, DOC_PROC_VECTOR_CELL_HEIGHT
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BACKGROUND_LIGHT
            cell.line.color.rgb = RGBColor(64, 64, 64)
            cell.line.width = Pt(1)

            # Cell text
            tf = cell.text_frame
            tf.text = value
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = FONT_SIZE_DOC_PROC_VECTOR
            p.font.color.rgb = COLOR_TEXT_GRAY
            for run in p.runs:
                run.font.name = FONT_FAMILY_MONOSPACE

        current_y += DOC_PROC_VECTOR_CELL_HEIGHT + Inches(0.2)

    # === BOTTOM: User Query Section ===
    # User icon (simple circle with "U")
    user_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2), DOC_PROC_QUERY_Y,
        DOC_PROC_USER_ICON_SIZE, DOC_PROC_USER_ICON_SIZE
    )
    user_circle.fill.background()
    user_circle.line.color.rgb = COLOR_ACCENT_BLUE
    user_circle.line.width = Pt(2)

    user_icon_text = slide.shapes.add_textbox(
        Inches(2), DOC_PROC_QUERY_Y,
        DOC_PROC_USER_ICON_SIZE, DOC_PROC_USER_ICON_SIZE
    )
    tf = user_icon_text.text_frame
    tf.text = "👤"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)

    # Query text
    query_box = slide.shapes.add_textbox(
        Inches(2.8), DOC_PROC_QUERY_Y,
        DOC_PROC_QUERY_TEXT_WIDTH, Inches(0.5)
    )
    tf = query_box.text_frame
    tf.text = '"What are the compliance requirements?"'
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = FONT_SIZE_DOC_PROC_QUERY
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.italic = True
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_REGULAR

    # Arrow
    arrow_box2 = slide.shapes.add_textbox(
        DOC_PROC_QUERY_ARROW_X, DOC_PROC_QUERY_Y,
        Inches(0.5), Inches(0.5)
    )
    tf = arrow_box2.text_frame
    tf.text = "→"
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_ACCENT_BLUE

    # Search term vector
    search_label = slide.shapes.add_textbox(
        DOC_PROC_SEARCH_X, DOC_PROC_QUERY_Y - Inches(0.3),
        Inches(2), Inches(0.25)
    )
    tf = search_label.text_frame
    tf.text = "Search term"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = FONT_SIZE_DOC_PROC_LABEL
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_BLUE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_SEMIBOLD

    # Search vector cells (highlighted in blue)
    search_values = ["0.19", "-0.73", "0.44", "0.88", "-0.31", "..."]
    cell_x = DOC_PROC_SEARCH_X
    cell_y = DOC_PROC_QUERY_Y + Inches(0.05)
    for i, value in enumerate(search_values):
        cell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cell_x + (i * (DOC_PROC_VECTOR_CELL_WIDTH + Inches(0.08))), cell_y,
            DOC_PROC_VECTOR_CELL_WIDTH, DOC_PROC_VECTOR_CELL_HEIGHT
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 60, 90)  # Dark blue background
        cell.line.color.rgb = COLOR_ACCENT_BLUE
        cell.line.width = Pt(2)

        # Cell text
        tf = cell.text_frame
        tf.text = value
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = FONT_SIZE_DOC_PROC_VECTOR
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_BLUE
        for run in p.runs:
            run.font.name = FONT_FAMILY_MONOSPACE

    return prs

def create_why_now_slide(prs, slide_num, step_data):
    """Helper function to create a single Why Now slide with one card"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_master_elements(slide, slide_num, 25)

    # Title - prominent at top
    title_box = slide.shapes.add_textbox(
        Inches(1), WHY_NOW_TITLE_Y,
        Inches(14), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.text = WHY_NOW_TITLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_WHY_NOW_TITLE
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_BLUE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_BOLD

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), WHY_NOW_SUBTITLE_Y,
        Inches(14), Inches(0.3)
    )
    tf = subtitle_box.text_frame
    tf.text = WHY_NOW_SUBTITLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_WHY_NOW_SUBTITLE
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_BLUE
    for run in p.runs:
        run.font.name = FONT_FAMILY_MONOSPACE

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        WHY_NOW_CARD_X, WHY_NOW_CARD_Y,
        WHY_NOW_CARD_WIDTH, WHY_NOW_CARD_HEIGHT
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(31, 41, 55)
    card.line.color.rgb = RGBColor(64, 64, 64)
    card.line.width = WHY_NOW_CARD_BORDER_WIDTH
    card.adjustments[0] = 0.08

    # Step number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        WHY_NOW_CARD_X + WHY_NOW_CIRCLE_X_OFFSET,
        WHY_NOW_CARD_Y + WHY_NOW_CIRCLE_Y_OFFSET,
        WHY_NOW_CIRCLE_SIZE, WHY_NOW_CIRCLE_SIZE
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    circle.line.width = Pt(0)

    # Number text
    num_box = slide.shapes.add_textbox(
        WHY_NOW_CARD_X + WHY_NOW_CIRCLE_X_OFFSET,
        WHY_NOW_CARD_Y + WHY_NOW_CIRCLE_Y_OFFSET,
        WHY_NOW_CIRCLE_SIZE, WHY_NOW_CIRCLE_SIZE
    )
    tf = num_box.text_frame
    tf.text = step_data["number"]
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_WHY_NOW_STEP_NUMBER
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_BOLD

    # Step title
    title_x = WHY_NOW_CARD_X + WHY_NOW_CONTENT_X_OFFSET
    title_y = WHY_NOW_CARD_Y + WHY_NOW_TITLE_Y_OFFSET
    step_title_box = slide.shapes.add_textbox(
        title_x, title_y,
        Inches(9), Inches(0.35)
    )
    tf = step_title_box.text_frame
    tf.text = step_data["title"]
    p = tf.paragraphs[0]
    p.font.size = FONT_SIZE_WHY_NOW_STEP_TITLE
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_SEMIBOLD

    # Bullets (stacked vertically)
    bullets_x = title_x
    bullets_y = WHY_NOW_CARD_Y + WHY_NOW_BULLETS_Y_OFFSET
    bullet_width = Inches(11)

    for i, (label, text) in enumerate(step_data["bullets"]):
        bullet_y = bullets_y + (i * (WHY_NOW_BULLET_HEIGHT + WHY_NOW_BULLET_GAP))

        bullet_box = slide.shapes.add_textbox(
            bullets_x, bullet_y,
            bullet_width, WHY_NOW_BULLET_HEIGHT
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True

        # Add text with bold label
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"▸ {label} "
        run1.font.size = FONT_SIZE_WHY_NOW_BULLET
        run1.font.bold = True
        run1.font.color.rgb = COLOR_TEXT_WHITE
        run1.font.name = FONT_FAMILY_INTER_SEMIBOLD

        run2 = p.add_run()
        run2.text = text
        run2.font.size = FONT_SIZE_WHY_NOW_BULLET
        run2.font.color.rgb = COLOR_TEXT_GRAY
        run2.font.name = FONT_FAMILY_INTER_REGULAR

    # Indicator badge
    indicator = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        WHY_NOW_CARD_X + WHY_NOW_INDICATOR_X_OFFSET,
        WHY_NOW_CARD_Y + WHY_NOW_INDICATOR_Y_OFFSET,
        WHY_NOW_INDICATOR_WIDTH, WHY_NOW_INDICATOR_HEIGHT
    )
    indicator.fill.background()
    indicator.line.color.rgb = step_data["indicator_color"]
    indicator.line.width = Pt(1)
    indicator.adjustments[0] = 0.25

    # Indicator text
    ind_text_box = slide.shapes.add_textbox(
        WHY_NOW_CARD_X + WHY_NOW_INDICATOR_X_OFFSET,
        WHY_NOW_CARD_Y + WHY_NOW_INDICATOR_Y_OFFSET,
        WHY_NOW_INDICATOR_WIDTH, WHY_NOW_INDICATOR_HEIGHT
    )
    tf = ind_text_box.text_frame
    tf.text = step_data["indicator"]
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = FONT_SIZE_WHY_NOW_INDICATOR
    p.font.bold = True
    p.font.color.rgb = step_data["indicator_color"]
    for run in p.runs:
        run.font.name = FONT_FAMILY_INTER_SEMIBOLD

    return prs

def create_slide_23(prs):
    """Slide 23: Why Now? - Card 1: AI Infrastructure Maturity"""
    step_data = {
        "number": "1",
        "title": "AI Infrastructure Maturity",
        "bullets": [
            ("Cost Revolution:", "Dramatic hardware price reductions alongside significant power efficiency improvements"),
            ("Container Management:", "Streamlined deployment and orchestration through Docker containerization technologies"),
            ("Performance & Memory:", "Remarkable computational performance gains with extensive shared memory capabilities"),
            ("Form & Operation:", "Ultra-compact form factors enabling whisper-quiet, enterprise-grade operation")
        ],
        "indicator": "⚡ TECHNICAL READINESS",
        "indicator_color": COLOR_WHY_NOW_TECH
    }
    return create_why_now_slide(prs, 23, step_data)

def create_slide_24(prs):
    """Slide 24: Why Now? - Card 2: Knowledge Worker Evolution"""
    step_data = {
        "number": "2",
        "title": "Knowledge Worker Evolution",
        "bullets": [
            ("Phase 1 - Playground:", "ChatGPT experimentation and individual productivity gains"),
            ("Phase 2 - Copilot:", "AI-assisted workflows and collaborative human-AI work"),
            ("Phase 3 - Workforce:", "Autonomous AI agents managing organizational knowledge"),
            ("Current Reality:", "73% of enterprises remain stuck between Phase 1-2")
        ],
        "indicator": "📈 MARKET DEMAND",
        "indicator_color": COLOR_WHY_NOW_MARKET
    }
    return create_why_now_slide(prs, 24, step_data)

def create_slide_25(prs):
    """Slide 25: Why Now? - Card 3: Data Sovereignty Crisis"""
    step_data = {
        "number": "3",
        "title": "Data Sovereignty Crisis",
        "bullets": [
            ("Regulatory Enforcement:", "EU AI Act and GDPR violations creating existential compliance risks"),
            ("Corporate Barriers:", "Fortune 500 companies cite data residency as primary AI adoption blocker"),
            ("Market Solution:", "Local plug-and-play AI becoming the gold standard for enterprises"),
            ("Value Proposition:", "Complete control, zero compliance risk, instant deployment capabilities")
        ],
        "indicator": "🏛️ REGULATORY FORCE",
        "indicator_color": COLOR_WHY_NOW_REGULATORY
    }
    return create_why_now_slide(prs, 25, step_data)

def create_presentation():
    """Creates the complete presentation with consistent master elements"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Slide 1: THE AI PARADOX
    create_slide_1(prs)

    # Slide 2: Organisations want AI
    create_slide_2(prs)

    # Slide 3: Market Reality
    create_slide_3(prs)

    # Slide 4: SOVEREIGN AI SOLUTION
    create_slide_4(prs)

    # Slide 5: TO BE DESIGNED
    create_slide_5(prs)

    # Slide 6: BRAIN-BRIDGES Hero Slide
    create_slide_6(prs)

    # Slide 7: TECHNICAL DEEP DIVE
    create_slide_7(prs)

    # Slide 8: Tokenization Intro - A Sample from legal domain
    create_slide_8(prs)

    # Slide 9: Vector Embeddings (Token → Vector Lookup)
    create_slide_9(prs)

    # Slide 10: Attention is all you need
    create_slide_10(prs)

    # Slide 11: Next word prediction
    create_slide_11(prs)

    # Slides 12-15: Autoregression (4 separate slides)
    create_slide_12(prs)  # Step 1
    create_slide_13(prs)  # Step 2
    create_slide_14(prs)  # Step 3
    create_slide_15(prs)  # Final

    # Slide 16: ON PREMISE MATTERS
    create_slide_16(prs)

    # Slide 17: The Fundamental Security Conflict
    create_slide_17(prs)

    # Slide 18: The Encryption Dilemma
    create_slide_18(prs)

    # Slide 19: Chat API Architecture
    create_slide_19(prs)

    # Slide 20: Chat API Architecture (copy)
    create_slide_20(prs)

    # Slide 21: RAG
    create_slide_21(prs)

    # Slide 22: Document Processing
    create_slide_22(prs)

    # Slide 23: Why Now? - Card 1: AI Infrastructure Maturity
    create_slide_23(prs)

    # Slide 24: Why Now? - Card 2: Knowledge Worker Evolution
    create_slide_24(prs)

    # Slide 25: Why Now? - Card 3: Data Sovereignty Crisis
    create_slide_25(prs)

    return prs

if __name__ == "__main__":
    print("🎨 Generating Brain-Bridges PowerPoint V3 with consistent master elements...")
    prs = create_presentation()

    # Create output directory if it doesn't exist
    os.makedirs("output", exist_ok=True)

    # Generate timestamp in format: YYYY_MM_DD___HH_MM_SS
    timestamp = datetime.now().strftime("%Y_%m_%d___%H_%M_%S")

    # Save timestamped version
    timestamped_path = f"output/{timestamp}__Brain-Bridges.pptx"
    prs.save(timestamped_path)
    print(f"✅ Timestamped version created: {timestamped_path}")

    # Save LATEST version (copy of timestamped file)
    latest_path = "output/Brain-Bridges_LATEST.pptx"
    shutil.copy2(timestamped_path, latest_path)
    print(f"✅ Latest version updated: {latest_path}")
    print("")
    print("📋 Slide Master Configuration:")
    print("   ✓ Background color: rgb(17, 24, 39)")
    print("   ✓ Logo 'BRAIN BRIDGES' top left (without v: xii)")
    print("   ✓ Slide counter top right")
    print("")
    print("🎯 Benefits:")
    print("   • New slides automatically inherit the design")
    print("   • Logo & slide numbers are always consistent")
    print("   • Master can be centrally adjusted")
    print("")
    print("💡 Tip: In PowerPoint under 'View' → 'Slide Master'")
    print("   you can edit the master!")
