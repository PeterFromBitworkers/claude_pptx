#!/usr/bin/env python3
"""
Test if python-pptx supports SVG files directly
"""

from pptx import Presentation
from pptx.util import Inches

try:
    # Create a test presentation
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Try to add an SVG
    svg_path = "assets/icons/legal.svg"
    picture = slide.shapes.add_picture(
        svg_path,
        Inches(1), Inches(1),
        width=Inches(1)
    )

    # Save test file
    prs.save("output/svg_test.pptx")
    print("✅ SUCCESS: SVG files are supported by python-pptx!")
    print(f"   Test file created: output/svg_test.pptx")
    print(f"   Picture type: {type(picture)}")

except Exception as e:
    print(f"❌ ERROR: SVG files are NOT directly supported")
    print(f"   Error message: {str(e)}")
    print(f"   Will need to convert to PNG instead")
